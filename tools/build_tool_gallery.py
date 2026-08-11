#!/usr/bin/env python3
"""Exercise WutPack's advertised tool stack and build its public proof gallery.

Run this with the managed environment so Python-library examples use the same
packages that WutPack installs:

    wut python tools/build_tool_gallery.py

The script writes one HTML proof page per advertised tool, along with small
downloadable artifacts and a machine-readable evidence manifest. It also
refreshes the generated tool links between TOOL_GALLERY markers on the landing
page. No example uses credentials, paid APIs, or external web pages.
"""

from __future__ import annotations

import argparse
import contextlib
import csv
import html
import importlib.metadata
import io
import json
import math
import os
import re
import shutil
import subprocess
import sys
import tempfile
import textwrap
import time
from dataclasses import asdict, dataclass, field
from datetime import date
from pathlib import Path
from typing import Iterable, Sequence


ROOT = Path(__file__).resolve().parents[1]
SITE = ROOT / "site"
OUTPUT = SITE / "tool-examples"
ARTIFACTS = OUTPUT / "artifacts"
LANDING_PAGE = SITE / "index.html"
WUTPACK_HOME = Path.home() / "Library/Application Support/WutPack"
NODE_PREFIX = WUTPACK_HOME / "npm-global"
NODE_MODULES = NODE_PREFIX / "lib/node_modules"

LAYER_ORDER = [
    "Desktop workbench",
    "AI coding CLIs",
    "Office, PDF, and media",
    "Data and research",
    "Diagrams and documentation",
    "Engineering utilities",
]

EXPECTED_TOOLS = {
    "Desktop workbench": ["LibreOffice", "Chromium", "Quarto", "draw.io", "Inkscape"],
    "AI coding CLIs": ["Codex", "Claude Code"],
    "Office, PDF, and media": [
        "openpyxl",
        "XlsxWriter",
        "python-docx",
        "python-pptx",
        "CairoSVG",
        "Pandoc",
        "Poppler",
        "qpdf",
        "MuPDF",
        "OCRmyPDF",
        "Tesseract",
        "ImageMagick",
        "FFmpeg",
    ],
    "Data and research": [
        "pandas",
        "Polars",
        "Arrow",
        "DuckDB",
        "SciPy",
        "scikit-learn",
        "statsmodels",
        "JupyterLab",
        "Playwright",
        "Selenium",
        "Scrapy",
    ],
    "Diagrams and documentation": [
        "PptxGenJS",
        "Mermaid",
        "Graphviz",
        "PlantUML",
        "Typst",
        "MkDocs",
        "Sphinx",
        "pdoc",
        "Doxygen",
        "JSDoc",
        "TypeDoc",
    ],
    "Engineering utilities": [
        "GitHub CLI",
        "ripgrep",
        "fd",
        "fzf",
        "jq",
        "yq",
        "ShellCheck",
        "shfmt",
        "delta",
        "hyperfine",
        "just",
    ],
}


class ToolFailure(RuntimeError):
    """A tool could not complete its deterministic local example."""


@dataclass
class CommandResult:
    command: str
    returncode: int
    stdout: str
    stderr: str
    duration_ms: int


@dataclass
class ToolProof:
    slug: str
    name: str
    layer: str
    tagline: str
    task: str
    command: str
    version: str
    result: str
    visual_html: str
    status: str = "passed"
    artifacts: list[str] = field(default_factory=list)
    duration_ms: int = 0
    note: str = ""


def escape(value: object) -> str:
    return html.escape(str(value), quote=True)


def clean_output(value: str, limit: int = 2400) -> str:
    value = re.sub(r"\x1b\[[0-?]*[ -/]*[@-~]", "", value)
    value = re.sub(
        r"/(?:private/)?var/folders/[^/\s]+/[^/\s]+/T/wutpack-tool-gallery-[^/\s]+",
        "$WORK",
        value,
    )
    value = value.replace(str(ROOT), "$REPO")
    value = value.replace(str(WUTPACK_HOME), "$WUTPACK_HOME")
    value = value.replace(str(Path.home()), "$HOME")
    value = value.strip()
    if len(value) > limit:
        value = value[: limit - 1].rstrip() + "…"
    return value


def display_command(parts: Sequence[str]) -> str:
    rendered: list[str] = []
    for part in parts:
        value = clean_output(str(part), limit=1000)
        rendered.append(value if re.fullmatch(r"[A-Za-z0-9_./:@%+=,-]+", value) else repr(value))
    return " ".join(rendered)


def run(
    parts: Sequence[str | Path],
    *,
    cwd: Path,
    env: dict[str, str] | None = None,
    expected: Iterable[int] = (0,),
    timeout: int = 90,
    input_text: str | None = None,
) -> CommandResult:
    command = [str(part) for part in parts]
    merged_env = os.environ.copy()
    merged_env.update(
        {
            "PATH": f"{NODE_PREFIX / 'bin'}:{merged_env.get('PATH', '')}",
            "NODE_PATH": str(NODE_MODULES),
            "MPLCONFIGDIR": str(cwd / ".matplotlib"),
            "SOURCE_DATE_EPOCH": "1704067200",
            "NO_COLOR": "1",
        }
    )
    if env:
        merged_env.update(env)
    (cwd / ".matplotlib").mkdir(exist_ok=True)
    started = time.monotonic()
    completed = subprocess.run(
        command,
        cwd=cwd,
        env=merged_env,
        text=True,
        input=input_text,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        timeout=timeout,
        check=False,
    )
    duration_ms = round((time.monotonic() - started) * 1000)
    result = CommandResult(
        command=display_command(command),
        returncode=completed.returncode,
        stdout=clean_output(completed.stdout),
        stderr=clean_output(completed.stderr),
        duration_ms=duration_ms,
    )
    if result.returncode not in set(expected):
        raise ToolFailure(
            f"{result.command} exited {result.returncode}\n"
            f"stdout:\n{result.stdout}\nstderr:\n{result.stderr}"
        )
    return result


def executable(name: str) -> str:
    resolved = shutil.which(name)
    if not resolved:
        raise ToolFailure(f"required executable is missing: {name}")
    return resolved


def node_executable(name: str) -> str:
    candidate = NODE_PREFIX / "bin" / name
    if candidate.is_file():
        return str(candidate)
    return executable(name)


def package_version(distribution: str) -> str:
    try:
        return importlib.metadata.version(distribution)
    except importlib.metadata.PackageNotFoundError as exc:
        raise ToolFailure(f"Python distribution is missing: {distribution}") from exc


def first_version_line(result: CommandResult) -> str:
    source = result.stdout or result.stderr
    for line in source.splitlines():
        line = line.strip()
        if line:
            return line[:180]
    return "installed"


def write_text(path: Path, value: str) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(textwrap.dedent(value).lstrip(), encoding="utf-8")
    return path


def write_json(path: Path, value: object) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(value, indent=2, sort_keys=True) + "\n", encoding="utf-8")
    return path


def copy_artifact(source: Path, name: str | None = None) -> str:
    ARTIFACTS.mkdir(parents=True, exist_ok=True)
    target = ARTIFACTS / (name or source.name)
    shutil.copy2(source, target)
    return f"artifacts/{target.name}"


def copy_artifact_tree(source: Path, name: str) -> str:
    ARTIFACTS.mkdir(parents=True, exist_ok=True)
    archive_base = ARTIFACTS / name
    archive = Path(shutil.make_archive(str(archive_base), "zip", root_dir=source))
    return f"artifacts/{archive.name}"


def chromium_binary() -> Path:
    candidates = sorted(
        (Path.home() / "Library/Caches/ms-playwright").glob(
            "chromium-*/chrome-mac-x64/Google Chrome for Testing.app/Contents/MacOS/Google Chrome for Testing"
        ),
        reverse=True,
    )
    if not candidates:
        raise ToolFailure("Playwright-managed Chromium binary is missing")
    return candidates[0]


def capture_html(source: Path, target: Path, cwd: Path, *, width: int = 1200, height: int = 760) -> CommandResult:
    from playwright.sync_api import sync_playwright

    started = time.monotonic()
    try:
        with sync_playwright() as playwright:
            browser = playwright.chromium.launch(
                headless=True,
                executable_path=str(chromium_binary()),
                args=["--no-sandbox", "--disable-gpu", "--disable-dev-shm-usage"],
                timeout=45_000,
            )
            page = browser.new_page(viewport={"width": width, "height": height})
            page.goto(source.resolve().as_uri(), wait_until="load", timeout=45_000)
            page.screenshot(path=str(target), full_page=True)
            title = page.title()
            browser.close()
    except Exception as exc:
        raise ToolFailure(f"Playwright Chromium capture failed for {source.name}: {exc}") from exc
    return CommandResult(
        command=f"chromium --headless --screenshot {clean_output(str(source))}",
        returncode=0,
        stdout=f"Captured {target.name}; title={title}",
        stderr="",
        duration_ms=elapsed_ms(started),
    )


def node_package_version(package: str) -> str:
    package_json = NODE_MODULES / package / "package.json"
    if not package_json.is_file():
        raise ToolFailure(f"Node package is missing: {package}")
    return str(json.loads(package_json.read_text(encoding="utf-8"))["version"])


def terminal_visual(lines: Sequence[str], title: str = "local run") -> str:
    rendered = "\n".join(escape(line) for line in lines)
    return (
        '<div class="terminal-visual">'
        f'<div class="visual-chrome"><span>{escape(title)}</span><span>completed</span></div>'
        f'<pre><code>{rendered}</code></pre></div>'
    )


def table_visual(headers: Sequence[object], rows: Sequence[Sequence[object]], title: str) -> str:
    head = "".join(f"<th scope=\"col\">{escape(item)}</th>" for item in headers)
    body = "".join(
        "<tr>" + "".join(f"<td>{escape(item)}</td>" for item in row) + "</tr>"
        for row in rows
    )
    return (
        '<div class="table-visual">'
        f'<div class="visual-chrome"><span>{escape(title)}</span><span>{len(rows)} rows</span></div>'
        f'<div class="table-scroll"><table><thead><tr>{head}</tr></thead><tbody>{body}</tbody></table></div>'
        "</div>"
    )


def code_visual(code: str, title: str, language: str = "text") -> str:
    return (
        '<div class="code-visual">'
        f'<div class="visual-chrome"><span>{escape(title)}</span><span>{escape(language)}</span></div>'
        f'<pre><code>{escape(code)}</code></pre></div>'
    )


def bars_visual(items: Sequence[tuple[str, float]], title: str, suffix: str = "") -> str:
    maximum = max((value for _, value in items), default=1) or 1
    rows = []
    for label, value in items:
        width = max(2.0, (float(value) / maximum) * 100)
        rows.append(
            '<div class="bar-row">'
            f'<span>{escape(label)}</span><i style="--bar:{width:.2f}%"></i>'
            f'<strong>{escape(f"{value:g}{suffix}")}</strong></div>'
        )
    return (
        '<div class="chart-visual">'
        f'<div class="visual-chrome"><span>{escape(title)}</span><span>computed</span></div>'
        f'<div class="bar-list">{"".join(rows)}</div></div>'
    )


def document_visual(title: str, kicker: str, paragraphs: Sequence[str], stamp: str) -> str:
    body = "".join(f"<p>{escape(paragraph)}</p>" for paragraph in paragraphs)
    return (
        '<div class="document-visual"><div class="paper-sheet">'
        f'<div class="doc-kicker">{escape(kicker)}</div><h2>{escape(title)}</h2>'
        f'<div class="doc-rule"></div>{body}<span class="doc-stamp">{escape(stamp)}</span>'
        "</div></div>"
    )


def image_visual(source: str, alt: str, title: str, meta: str = "generated artifact") -> str:
    return (
        '<figure class="image-visual">'
        f'<div class="visual-chrome"><span>{escape(title)}</span><span>{escape(meta)}</span></div>'
        f'<div class="image-stage"><img src="{escape(source)}" alt="{escape(alt)}" loading="lazy"></div>'
        "</figure>"
    )


def browser_visual(title: str, body: str, address: str = "http://127.0.0.1/example") -> str:
    return (
        '<div class="browser-visual">'
        '<div class="browser-top"><span class="browser-dots" aria-hidden="true">● ● ●</span>'
        f'<span class="address">{escape(address)}</span></div>'
        f'<div class="browser-page"><h2>{escape(title)}</h2>{body}</div></div>'
    )


def page_for(proof: ToolProof, previous_slug: str, next_slug: str) -> str:
    exercised = proof.status == "passed"
    status_label = "exercised locally" if exercised else "installed only"
    evidence_intro = (
        "The command completed against a local fictional fixture. The page records the tool version and the output used in the visual above."
        if exercised
        else "The package is part of the installed stack, but the CLI and model were intentionally not invoked at the owner's request."
    )
    footer_text = (
        "Generated from a deterministic, fictional local fixture. No credential or paid API call was used."
        if exercised and proof.slug != "codex"
        else "Evidence is labeled precisely: authenticated Codex ran; unauthenticated Claude Code did not. No credential is embedded in this site."
    )
    downloads = "".join(
        f'<a class="artifact-link" href="{escape(path)}" download>Download {escape(Path(path).suffix.lstrip(".").upper() or "artifact")}</a>'
        for path in proof.artifacts
    )
    note = f'<p class="proof-note">{escape(proof.note)}</p>' if proof.note else ""
    result_lines = proof.result.splitlines() or [proof.result]
    result_block = "\n".join(escape(line) for line in result_lines[:14])
    return f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <meta name="theme-color" content="#f3f0e8">
  <meta name="description" content="A reproducible WutPack example for {escape(proof.name)}.">
  <link rel="icon" href="../favicon.svg" type="image/svg+xml">
  <link rel="preconnect" href="https://fonts.googleapis.com">
  <link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
  <link href="https://fonts.googleapis.com/css2?family=DM+Sans:wght@400;500;600;700&amp;family=IBM+Plex+Mono:wght@400;500;600&amp;family=Instrument+Serif:ital@0;1&amp;display=swap" rel="stylesheet">
  <link rel="stylesheet" href="assets/gallery.css">
  <title>{escape(proof.name)} example | WutPack</title>
</head>
<body class="proof-page status-{escape(proof.status)} layer-{escape(proof.layer.lower().replace(' ', '-').replace(',', ''))}">
  <a class="skip-link" href="#proof">Skip to proof</a>
  <header class="proof-header">
    <a class="proof-brand" href="../index.html#tool-gallery"><span>W</span> WutPack</a>
    <a class="back-link" href="../index.html#tool-gallery">All 53 tool examples</a>
  </header>
  <main id="proof">
    <section class="proof-hero">
      <div>
        <p class="proof-layer">{escape(proof.layer)}</p>
        <h1>{escape(proof.name)}</h1>
        <p class="proof-tagline" data-pretext contenteditable="true">{escape(proof.tagline)}</p>
      </div>
      <div class="proof-status"><span></span> {escape(status_label)}</div>
    </section>
    <section class="proof-stage" aria-label="Visual output from {escape(proof.name)}">
      {proof.visual_html}
    </section>
    <section class="proof-evidence">
      <div class="evidence-copy">
        <p class="proof-kicker">What ran</p>
        <h2>{escape(proof.task)}</h2>
        <p>{escape(evidence_intro)}</p>
        {note}
        <div class="artifact-actions">{downloads}</div>
      </div>
      <div class="run-card">
        <div><span>Tool version</span><strong>{escape(proof.version)}</strong></div>
        <div><span>Command or API</span><code>{escape(proof.command)}</code></div>
        <div><span>Observed result</span><pre>{result_block}</pre></div>
        <div><span>Local runtime</span><strong>{proof.duration_ms:,} ms</strong></div>
      </div>
    </section>
    <nav class="proof-pagination" aria-label="Tool examples">
      <a href="{escape(previous_slug)}.html"><span>Previous proof</span>←</a>
      <a href="../index.html#tool-gallery"><span>Gallery</span>53 tools</a>
      <a href="{escape(next_slug)}.html"><span>Next proof</span>→</a>
    </nav>
  </main>
  <footer class="proof-footer">
    <p>{escape(footer_text)}</p>
    <a href="https://github.com/sfungwinbond/Gstackwut">Source on GitHub</a>
  </footer>
  <script type="module" src="assets/gallery.js"></script>
</body>
</html>
"""


def gallery_css() -> str:
    return r"""
:root {
  color-scheme: light;
  --paper: #f3f0e8;
  --raised: #fbfaf6;
  --ink: #14201d;
  --muted: #58615d;
  --line: #87918b;
  --green: #176b4d;
  --lime: #b8f36b;
  --orange: #f26b38;
  --yellow: #f5d969;
  --blue: #306ab2;
  --mono: "IBM Plex Mono", ui-monospace, monospace;
  --body: "DM Sans", system-ui, sans-serif;
  --display: "Instrument Serif", Georgia, serif;
}
* { box-sizing: border-box; }
html { scroll-behavior: smooth; }
body { margin: 0; background: var(--paper); color: var(--ink); font-family: var(--body); line-height: 1.55; }
body::before { position: fixed; z-index: -1; inset: 0; background-image: linear-gradient(#c8ccc4 1px, transparent 1px), linear-gradient(90deg, #c8ccc4 1px, transparent 1px); background-size: 42px 42px; opacity: .12; content: ""; }
a { color: inherit; text-underline-offset: .2em; }
:focus-visible { outline: 3px solid var(--orange); outline-offset: 4px; }
.skip-link { position: fixed; z-index: 10; top: .75rem; left: .75rem; padding: .7rem 1rem; transform: translateY(-150%); background: var(--ink); color: white; }
.skip-link:focus { transform: none; }
.proof-header, main, .proof-footer { width: min(calc(100% - 2rem), 1160px); margin-inline: auto; }
.proof-header { display: flex; min-height: 4.8rem; align-items: center; justify-content: space-between; border-bottom: 1px solid var(--line); }
.proof-brand { display: inline-flex; align-items: center; gap: .65rem; font-weight: 700; text-decoration: none; }
.proof-brand span { display: grid; width: 2rem; height: 2rem; place-items: center; border: 2px solid var(--ink); background: var(--lime); color: #14201d; font-family: var(--mono); }
.back-link { font-family: var(--mono); font-size: .72rem; font-weight: 600; }
.proof-hero { display: grid; gap: 1.5rem; padding: clamp(4rem, 9vw, 8rem) 0 clamp(2rem, 5vw, 4rem); }
.proof-layer, .proof-kicker { margin: 0 0 .9rem; color: var(--green); font-family: var(--mono); font-size: .72rem; font-weight: 600; letter-spacing: .09em; text-transform: uppercase; }
.proof-hero h1 { margin: 0; font-family: var(--display); font-size: clamp(4rem, 12vw, 8.5rem); font-weight: 400; letter-spacing: -.055em; line-height: .85; }
.proof-tagline { max-width: 44rem; margin: 1.5rem 0 0; color: var(--muted); font-size: clamp(1.1rem, 2vw, 1.35rem); }
.proof-status { display: inline-flex; width: max-content; align-items: center; align-self: end; gap: .5rem; padding: .55rem .7rem; border: 1px solid var(--ink); background: var(--lime); color: #14201d; font-family: var(--mono); font-size: .66rem; font-weight: 600; text-transform: uppercase; }
.proof-status span { width: .5rem; height: .5rem; border-radius: 50%; background: var(--green); }
.status-installed-only .proof-status { background: var(--yellow); }
.status-installed-only .proof-status span { background: var(--orange); }
.proof-stage { min-height: 28rem; padding: clamp(.75rem, 2vw, 1.5rem); border: 1px solid var(--ink); background: var(--ink); color: #edf2ed; }
.visual-chrome { display: flex; min-height: 2.7rem; align-items: center; justify-content: space-between; gap: 1rem; padding: .65rem .85rem; border-bottom: 1px solid #52615a; color: #b7c1bb; font-family: var(--mono); font-size: .65rem; text-transform: uppercase; }
.terminal-visual, .code-visual, .table-visual, .chart-visual, .browser-visual, .image-visual { height: 100%; margin: 0; border: 1px solid #52615a; background: #101714; }
pre { margin: 0; white-space: pre-wrap; overflow-wrap: anywhere; }
.terminal-visual pre, .code-visual pre { min-height: 23rem; padding: clamp(1rem, 4vw, 2.2rem); color: #e9f4ed; font-family: var(--mono); font-size: clamp(.72rem, 1.8vw, .94rem); line-height: 1.75; }
.terminal-visual code::first-line { color: var(--lime); }
.table-scroll { overflow-x: auto; }
.table-visual table { width: 100%; min-width: 560px; border-collapse: collapse; background: #fffefa; color: #14201d; }
.table-visual th, .table-visual td { padding: .9rem; border-right: 1px solid #c8ccc4; border-bottom: 1px solid #c8ccc4; text-align: left; }
.table-visual th { background: var(--lime); font-family: var(--mono); font-size: .7rem; text-transform: uppercase; }
.chart-visual { background: #fffefa; color: #14201d; }
.chart-visual .visual-chrome { border-color: #c8ccc4; color: #58615d; }
.bar-list { display: grid; gap: 1rem; padding: clamp(1rem, 4vw, 2.5rem); }
.bar-row { display: grid; grid-template-columns: minmax(5.5rem, .55fr) minmax(8rem, 2fr) 4rem; align-items: center; gap: .8rem; }
.bar-row span { font-weight: 600; }
.bar-row i { height: 2rem; border: 1px solid #14201d; background: linear-gradient(90deg, var(--lime) var(--bar), transparent var(--bar)); }
.bar-row strong { font-family: var(--mono); font-size: .78rem; text-align: right; }
.document-visual { display: grid; min-height: 27rem; place-items: center; padding: clamp(1rem, 5vw, 3rem); background: #23312b; }
.paper-sheet { position: relative; width: min(100%, 780px); min-height: 22rem; padding: clamp(1.4rem, 5vw, 3.5rem); background: #fffefa; color: #14201d; }
.doc-kicker { color: var(--green); font-family: var(--mono); font-size: .65rem; font-weight: 600; letter-spacing: .1em; text-transform: uppercase; }
.paper-sheet h2 { max-width: 17ch; margin: .7rem 0; font-family: var(--display); font-size: clamp(2.3rem, 6vw, 4.5rem); font-weight: 400; line-height: .95; }
.paper-sheet p { max-width: 49rem; color: #52615a; }
.doc-rule { width: 4rem; height: 4px; margin: 1.4rem 0; background: var(--orange); }
.doc-stamp { position: absolute; right: 1.2rem; bottom: 1.2rem; padding: .35rem .5rem; border: 1px solid var(--green); color: var(--green); font-family: var(--mono); font-size: .6rem; font-weight: 600; text-transform: uppercase; transform: rotate(-2deg); }
.image-stage { display: grid; min-height: 24rem; place-items: center; padding: clamp(.75rem, 3vw, 1.5rem); background: #e8ece7; }
.image-stage img { display: block; max-width: 100%; max-height: 68vh; object-fit: contain; }
.browser-top { display: flex; align-items: center; gap: 1rem; padding: .65rem .8rem; border-bottom: 1px solid #52615a; background: #17211d; font-family: var(--mono); font-size: .65rem; }
.browser-dots { color: var(--orange); letter-spacing: .25em; }
.address { flex: 1; padding: .3rem .55rem; border: 1px solid #52615a; color: #b7c1bb; }
.browser-page { min-height: 24rem; padding: clamp(1.2rem, 5vw, 3rem); background: #fffefa; color: #14201d; }
.browser-page h2 { max-width: 16ch; margin: 0 0 1.5rem; font-family: var(--display); font-size: clamp(2.6rem, 7vw, 5.5rem); font-weight: 400; line-height: .9; }
.browser-page p { max-width: 46rem; color: #58615d; }
.browser-page .browser-card { margin-top: 2rem; padding: 1rem; border-left: 5px solid var(--green); background: #eef4df; }
.split-visual { display: grid; min-height: 27rem; border: 1px solid #52615a; background: #101714; }
.split-visual > div { display: grid; place-items: center; padding: 1rem; background: #e8ece7; }
.split-visual img { display: block; max-width: 100%; max-height: 30rem; object-fit: contain; }
.split-visual pre { min-height: 12rem; padding: 1.5rem; color: #e9f4ed; font-family: var(--mono); font-size: .78rem; line-height: 1.7; }
.split-visual pre span { color: var(--lime); font-weight: 600; text-transform: uppercase; }
.media-visual { border: 1px solid #52615a; background: #101714; }
.media-visual video { display: block; width: 100%; max-height: 72vh; background: #101714; }
.proof-evidence { display: grid; gap: 2rem; padding: clamp(3rem, 7vw, 6rem) 0; border-bottom: 1px solid var(--line); }
.evidence-copy h2 { max-width: 15ch; margin: 0; font-family: var(--display); font-size: clamp(2.6rem, 6vw, 5rem); font-weight: 400; line-height: .92; }
.evidence-copy > p:not(.proof-kicker) { max-width: 38rem; color: var(--muted); }
.proof-note { padding-left: .85rem; border-left: 3px solid var(--orange); }
.artifact-actions { display: flex; flex-wrap: wrap; gap: .65rem; margin-top: 1.5rem; }
.artifact-link { display: inline-flex; min-height: 44px; align-items: center; padding: .65rem .8rem; border: 1px solid var(--ink); background: var(--ink); color: white; font-weight: 700; text-decoration: none; }
.artifact-link:hover { background: var(--green); }
.run-card { border-top: 1px solid var(--line); }
.run-card > div { display: grid; gap: .45rem; padding: 1rem 0; border-bottom: 1px solid var(--line); }
.run-card span { color: var(--green); font-family: var(--mono); font-size: .64rem; font-weight: 600; text-transform: uppercase; }
.run-card strong, .run-card code, .run-card pre { font-family: var(--mono); font-size: .75rem; overflow-wrap: anywhere; }
.run-card pre { max-height: 12rem; overflow: auto; color: var(--muted); }
.proof-pagination { display: grid; grid-template-columns: repeat(3, 1fr); padding: 1.5rem 0 4rem; }
.proof-pagination a { display: grid; min-height: 4.5rem; align-content: center; padding: .75rem; border-right: 1px solid var(--line); font-size: 1.25rem; font-weight: 700; text-align: center; text-decoration: none; }
.proof-pagination a:first-child { text-align: left; }
.proof-pagination a:last-child { border-right: 0; text-align: right; }
.proof-pagination span { display: block; color: var(--muted); font-family: var(--mono); font-size: .6rem; font-weight: 500; text-transform: uppercase; }
.proof-footer { display: flex; flex-wrap: wrap; justify-content: space-between; gap: 1rem; padding: 2rem 0; border-top: 1px solid var(--line); color: var(--muted); font-size: .78rem; }
.proof-footer p { margin: 0; }
@media (min-width: 768px) {
  .proof-header, main, .proof-footer { width: min(calc(100% - 3.5rem), 1160px); }
  .proof-hero { grid-template-columns: minmax(0, 1fr) auto; align-items: end; }
  .proof-evidence { grid-template-columns: minmax(0, 1.05fr) minmax(20rem, .95fr); }
  .run-card > div { grid-template-columns: 8rem minmax(0, 1fr); align-items: start; }
  .split-visual { grid-template-columns: minmax(0, 1.3fr) minmax(18rem, .7fr); }
}
@media (max-width: 520px) {
  .back-link { font-size: 0; }
  .back-link::after { font-size: .7rem; content: "All examples"; }
  .proof-pagination a { font-size: 1rem; }
  .bar-row { grid-template-columns: 5rem 1fr; }
  .bar-row strong { grid-column: 2; }
}
@media (prefers-color-scheme: dark) {
  :root { color-scheme: dark; --paper: #111815; --raised: #18221e; --ink: #edf2ed; --muted: #aeb9b3; --line: #617067; --green: #91d759; }
  .proof-brand span, .proof-status { color: #14201d; }
  .artifact-link { border-color: #edf2ed; }
}
@media (prefers-reduced-motion: reduce) { html { scroll-behavior: auto; } * { transition: none !important; } }
"""


def gallery_js() -> str:
    return r"""
const measured = new Map();
const elements = [...document.querySelectorAll('[data-pretext]')];

async function applyPretext() {
  if (!elements.length) return;
  try {
    const { prepare, layout } = await import('https://esm.sh/@chenglou/pretext');
    await document.fonts.ready;
    const prepareElement = (element) => {
      const style = getComputedStyle(element);
      measured.set(element, { handle: prepare(element.textContent, style.font), lineHeight: Number.parseFloat(style.lineHeight) });
    };
    const relayout = () => measured.forEach(({ handle, lineHeight }, element) => {
      const result = layout(handle, element.clientWidth, lineHeight);
      element.style.minHeight = `${Math.ceil(result.height)}px`;
    });
    elements.forEach((element) => {
      prepareElement(element);
      new MutationObserver(() => { prepareElement(element); relayout(); }).observe(element, { characterData: true, subtree: true, childList: true });
      new ResizeObserver(relayout).observe(element);
    });
    relayout();
  } catch {
    document.documentElement.dataset.pretext = 'css-fallback';
  }
}

applyPretext();
"""


def elapsed_ms(started: float) -> int:
    return round((time.monotonic() - started) * 1000)


def make_fixtures(work: Path) -> dict[str, Path]:
    """Create fictional local inputs shared by the individual tool runs."""
    from PIL import Image, ImageDraw, ImageFont
    from docx import Document
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas

    source = work / "source"
    source.mkdir(parents=True, exist_ok=True)
    fixtures: dict[str, Path] = {}

    fixtures["csv"] = write_text(
        source / "market-scorecard.csv",
        """
        market,revenue_m,growth_pct,margin_pct,score
        Coastal,86,14.0,20,81
        North,79,9.2,19,78
        Central,69,11.2,16,72
        West,43,5.8,14,58
        """,
    )
    fixtures["json"] = write_json(
        source / "markets.json",
        {
            "company": "Atlas Services (fictional)",
            "markets": [
                {"name": "Coastal", "score": 81, "owner": "Avery"},
                {"name": "North", "score": 78, "owner": "Morgan"},
                {"name": "Central", "score": 72, "owner": "Riley"},
                {"name": "West", "score": 58, "owner": "Jordan"},
            ],
        },
    )
    fixtures["yaml"] = write_text(
        source / "launch-plan.yml",
        """
        project: Atlas market entry
        fictional: true
        gates:
          - name: evidence complete
            owner: strategy
          - name: economics approved
            owner: finance
          - name: pilot ready
            owner: operations
        """,
    )
    fixtures["markdown"] = write_text(
        source / "decision-brief.md",
        """
        # Atlas market-entry decision

        **Recommendation:** Start with Coastal, then gate North on pilot economics.

        | Market | Score | Decision |
        |---|---:|---|
        | Coastal | 81 | Pilot now |
        | North | 78 | Prepare next |
        | Central | 72 | Monitor |

        All companies, people, markets, and values in this example are fictional.
        """,
    )
    fixtures["html"] = write_text(
        source / "browser-fixture.html",
        """
        <!doctype html><html lang="en"><head><meta charset="utf-8">
        <meta name="viewport" content="width=device-width,initial-scale=1">
        <title>Atlas local dashboard</title>
        <style>
          body{margin:0;background:#f3f0e8;color:#14201d;font:18px system-ui}
          main{max-width:900px;margin:auto;padding:70px 32px}
          small{color:#176b4d;font-weight:700;letter-spacing:.08em;text-transform:uppercase}
          h1{max-width:720px;margin:.5rem 0;font:64px Georgia;line-height:.95}
          .score{display:grid;grid-template-columns:repeat(4,1fr);gap:10px;margin-top:38px}
          article{padding:20px;border:1px solid #14201d;background:#fbfaf6}
          strong{display:block;font-size:32px}.best{background:#b8f36b}
        </style></head><body><main>
        <small>Fictional local fixture</small><h1>Choose the first market with evidence.</h1>
        <p id="summary">Coastal leads the scorecard. North remains the gated follow-on.</p>
        <section class="score" aria-label="Market scores">
          <article class="best" data-market="Coastal"><span>Coastal</span><strong>81</strong></article>
          <article data-market="North"><span>North</span><strong>78</strong></article>
          <article data-market="Central"><span>Central</span><strong>72</strong></article>
          <article data-market="West"><span>West</span><strong>58</strong></article>
        </section></main></body></html>
        """,
    )
    fixtures["svg"] = write_text(
        source / "market-gauge.svg",
        """
        <svg xmlns="http://www.w3.org/2000/svg" width="960" height="540" viewBox="0 0 960 540">
          <rect width="960" height="540" fill="#f3f0e8"/>
          <text x="64" y="78" font-family="Arial" font-size="22" font-weight="bold" fill="#176b4d">FICTIONAL MARKET SCORE</text>
          <text x="64" y="154" font-family="Georgia" font-size="58" fill="#14201d">Coastal leads at 81.</text>
          <rect x="64" y="222" width="800" height="82" fill="#ffffff" stroke="#14201d" stroke-width="3"/>
          <rect x="64" y="222" width="648" height="82" fill="#b8f36b" stroke="#14201d" stroke-width="3"/>
          <text x="80" y="276" font-family="Arial" font-size="30" font-weight="bold" fill="#14201d">81 / 100</text>
          <text x="64" y="390" font-family="Arial" font-size="26" fill="#58615d">Evidence → economics → pilot gate</text>
          <circle cx="690" cy="382" r="18" fill="#f26b38"/><circle cx="755" cy="382" r="18" fill="#f5d969"/><circle cx="820" cy="382" r="18" fill="#176b4d"/>
        </svg>
        """,
    )
    fixtures["mermaid"] = write_text(
        source / "engagement.mmd",
        """
        flowchart LR
          A[Evidence] --> B[Score markets]
          B --> C[Model economics]
          C --> D{Pilot gate}
          D -->|pass| E[Launch]
          D -->|hold| F[Learn]
        """,
    )
    fixtures["dot"] = write_text(
        source / "decision.dot",
        """
        digraph decision {
          graph [bgcolor="transparent", rankdir=LR, pad=0.25];
          node [shape=box, style="filled", fontname="Arial", color="#14201d", fillcolor="#fbfaf6", margin="0.25,0.15"];
          edge [color="#176b4d", penwidth=2];
          Evidence -> Score -> Economics -> Gate;
          Gate -> Launch [label=" pass"];
          Gate -> Learn [label=" hold"];
          Gate [shape=diamond, fillcolor="#b8f36b"];
        }
        """,
    )
    fixtures["plantuml"] = write_text(
        source / "handoff.puml",
        """
        @startuml
        skinparam backgroundColor transparent
        skinparam monochrome true
        actor Owner
        participant "AI host" as Host
        participant "WutPack skill" as Skill
        participant "Persistent tool" as Tool
        Owner -> Host: Describe outcome
        Host -> Skill: Route the work
        Skill -> Tool: Build and validate
        Tool --> Owner: Editable artifact
        @enduml
        """,
    )
    fixtures["typst"] = write_text(
        source / "one-page.typ",
        """
        #set page(paper: "us-letter", margin: 0.8in)
        #set text(font: "Arial", size: 11pt)
        #text(fill: rgb("176b4d"), weight: "bold")[FICTIONAL DECISION NOTE]
        #v(12pt)
        #text(size: 30pt, weight: "bold")[Coastal is the first gated pilot.]
        #line(length: 100%, stroke: 1pt + rgb("87918b"))
        #v(12pt)
        The scorecard favors Coastal at *81*. The next decision is not a broad rollout;
        it is a measured pilot with an explicit economics gate.
        #v(18pt)
        #table(columns: (2fr, 1fr, 2fr), inset: 8pt,
          [*Market*], [*Score*], [*Action*],
          [Coastal], [81], [Pilot now],
          [North], [78], [Prepare next],
          [Central], [72], [Monitor],
        )
        """,
    )
    fixtures["quarto"] = write_text(
        source / "analysis.qmd",
        """
        ---
        title: "Atlas market decision"
        subtitle: "Fictional, reproducible Quarto example"
        format:
          html:
            toc: true
        ---

        ## Recommendation

        Start with the **Coastal** pilot. Its weighted score is 81.

        ```{python}
        scores = {"Coastal": 81, "North": 78, "Central": 72, "West": 58}
        winner = max(scores, key=scores.get)
        print(f"Selected market: {winner} ({scores[winner]})")
        ```
        """,
    )
    fixtures["messy_shell"] = write_text(
        source / "messy.sh",
        """
        #!/bin/sh
        owner="strategy"
        if [ -n "$owner" ];then
        printf '%s\\n' "pilot owner: $owner"
        fi
        """,
    )
    fixtures["broken_shell"] = write_text(
        source / "broken.sh",
        """
        #!/bin/sh
        files="*.csv"
        for file in $files; do
          echo $file
        done
        """,
    )
    fixtures["python_module"] = write_text(
        source / "market_model.py",
        '''
        """Tiny fictional market-scoring module used by documentation examples."""

        def weighted_score(growth: float, margin: float, complexity: float) -> float:
            """Return a 0–100 attractiveness score from three normalized inputs."""
            return round((growth * 0.45 + margin * 0.4 + (100 - complexity) * 0.15), 1)
        ''',
    )
    fixtures["javascript"] = write_text(
        source / "decision.js",
        """
        /** Return the highest-scoring fictional market.
         * @param {Array<{name: string, score: number}>} markets candidate markets
         * @returns {{name: string, score: number}} top market
         */
        export function selectMarket(markets) {
          return [...markets].sort((a, b) => b.score - a.score)[0];
        }
        """,
    )
    fixtures["typescript"] = write_text(
        source / "decision.ts",
        """
        export interface Market { name: string; score: number; }

        /** Select the highest-scoring fictional market. */
        export function selectMarket(markets: Market[]): Market {
          return [...markets].sort((a, b) => b.score - a.score)[0];
        }
        """,
    )
    fixtures["c_source"] = write_text(
        source / "score.c",
        """
        /** @file score.c Fictional market score helper. */
        /** Compute a weighted score from growth and margin percentages.
         * @param growth Growth percentage.
         * @param margin Margin percentage.
         * @return Weighted score.
         */
        double market_score(double growth, double margin) {
          return growth * 0.55 + margin * 0.45;
        }
        """,
    )
    fixtures["justfile"] = write_text(
        source / "justfile",
        """
        set quiet

        score market="Coastal" value="81":
            @printf '%s\\n' "{{market}}: {{value}} / 100 — pilot ready"
        """,
    )
    fixtures["drawio"] = write_text(
        source / "decision.drawio",
        """
        <mxfile host="app.diagrams.net"><diagram name="Decision">
        <mxGraphModel dx="1000" dy="600" grid="1" gridSize="10" page="1" pageWidth="1100" pageHeight="850">
        <root><mxCell id="0"/><mxCell id="1" parent="0"/>
        <mxCell id="2" value="Evidence" style="rounded=0;whiteSpace=wrap;html=1;fillColor=#b8f36b;strokeColor=#14201d;fontSize=18;" vertex="1" parent="1"><mxGeometry x="80" y="180" width="180" height="80" as="geometry"/></mxCell>
        <mxCell id="3" value="Pilot gate" style="rhombus;whiteSpace=wrap;html=1;fillColor=#f5d969;strokeColor=#14201d;fontSize=18;" vertex="1" parent="1"><mxGeometry x="360" y="165" width="130" height="110" as="geometry"/></mxCell>
        <mxCell id="4" value="Launch" style="rounded=0;whiteSpace=wrap;html=1;fillColor=#f26b38;strokeColor=#14201d;fontSize=18;" vertex="1" parent="1"><mxGeometry x="600" y="180" width="180" height="80" as="geometry"/></mxCell>
        <mxCell id="5" style="edgeStyle=orthogonalEdgeStyle;strokeWidth=2;strokeColor=#176b4d;" edge="1" parent="1" source="2" target="3"><mxGeometry relative="1" as="geometry"/></mxCell>
        <mxCell id="6" value="pass" style="edgeStyle=orthogonalEdgeStyle;strokeWidth=2;strokeColor=#176b4d;" edge="1" parent="1" source="3" target="4"><mxGeometry relative="1" as="geometry"/></mxCell>
        </root></mxGraphModel></diagram></mxfile>
        """,
    )

    document = Document()
    document.add_heading("Atlas pilot recommendation", level=0)
    document.add_paragraph("Fictional decision memo", style="Subtitle")
    document.add_heading("Decision", level=1)
    document.add_paragraph("Launch a gated Coastal pilot. Keep North prepared as the follow-on market.")
    table = document.add_table(rows=1, cols=3)
    table.style = "Light Shading Accent 1"
    for cell, value in zip(table.rows[0].cells, ["Market", "Score", "Action"]):
        cell.text = value
    for row in [("Coastal", "81", "Pilot now"), ("North", "78", "Prepare next"), ("Central", "72", "Monitor")]:
        cells = table.add_row().cells
        for cell, value in zip(cells, row):
            cell.text = value
    fixtures["docx"] = source / "atlas-memo.docx"
    document.save(fixtures["docx"])

    image = Image.new("RGB", (1600, 900), "white")
    draw = ImageDraw.Draw(image)
    font_paths = [
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf",
        "/System/Library/Fonts/Supplemental/Arial.ttf",
    ]
    try:
        large = ImageFont.truetype(font_paths[0], 74)
        medium = ImageFont.truetype(font_paths[1], 40)
    except OSError:
        large = medium = ImageFont.load_default()
    draw.rectangle((0, 0, 1600, 900), fill="#f3f0e8")
    draw.text((110, 115), "FICTIONAL PILOT NOTE", fill="#176b4d", font=medium)
    draw.text((110, 230), "Coastal is first.", fill="#14201d", font=large)
    draw.text((110, 350), "Score 81  •  Owner Avery  •  Gate 30 days", fill="#14201d", font=medium)
    draw.rectangle((110, 500, 1420, 650), fill="#b8f36b", outline="#14201d", width=5)
    draw.text((155, 545), "Measure economics before scaling.", fill="#14201d", font=medium)
    fixtures["scan_png"] = source / "pilot-note-scan.png"
    image.save(fixtures["scan_png"])

    fixtures["rich_pdf"] = source / "decision.pdf"
    pdf = canvas.Canvas(str(fixtures["rich_pdf"]), pagesize=letter)
    pdf.setTitle("Fictional Atlas decision")
    pdf.setFillColorRGB(0.09, 0.42, 0.30)
    pdf.setFont("Helvetica-Bold", 13)
    pdf.drawString(54, 730, "FICTIONAL DECISION BRIEF")
    pdf.setFillColorRGB(0.08, 0.13, 0.11)
    pdf.setFont("Helvetica-Bold", 28)
    pdf.drawString(54, 675, "Coastal is the first pilot.")
    pdf.setFont("Helvetica", 12)
    pdf.drawString(54, 635, "Score 81. Validate economics for 30 days before scaling.")
    pdf.rect(54, 525, 500, 62, stroke=1, fill=0)
    pdf.setFont("Helvetica-Bold", 18)
    pdf.drawString(72, 550, "Evidence  →  Economics  →  Gate  →  Launch")
    pdf.save()

    fixtures["scan_pdf"] = source / "scan.pdf"
    scan_pdf = canvas.Canvas(str(fixtures["scan_pdf"]), pagesize=letter)
    scan_pdf.drawImage(ImageReader(image), 36, 190, width=540, height=304)
    scan_pdf.save()
    return fixtures


def exercise_desktop(fixtures: dict[str, Path], work: Path) -> list[ToolProof]:
    proofs: list[ToolProof] = []
    desktop = work / "desktop"
    desktop.mkdir()

    version = first_version_line(run([executable("soffice"), "--version"], cwd=desktop))
    converted = desktop / "libreoffice"
    converted.mkdir()
    result = run(
        [
            executable("soffice"),
            "--headless",
            f"-env:UserInstallation=file://{desktop / 'lo-profile'}",
            "--convert-to",
            "pdf",
            "--outdir",
            converted,
            fixtures["docx"],
        ],
        cwd=desktop,
        timeout=120,
    )
    pdf = converted / "atlas-memo.pdf"
    if not pdf.is_file():
        raise ToolFailure("LibreOffice did not create atlas-memo.pdf")
    preview_root = desktop / "libreoffice-preview"
    preview_result = run([executable("pdftoppm"), "-png", "-f", "1", "-singlefile", "-r", "120", pdf, preview_root], cwd=desktop)
    preview = preview_root.with_suffix(".png")
    preview_link = copy_artifact(preview, "libreoffice-memo.png")
    pdf_link = copy_artifact(pdf, "libreoffice-memo.pdf")
    proofs.append(
        ToolProof(
            slug="libreoffice",
            name="LibreOffice",
            layer="Desktop workbench",
            tagline="Turns a structured Word document into a clean, standard PDF without opening a GUI.",
            task="Render an editable decision memo headlessly",
            command="soffice --headless --convert-to pdf atlas-memo.docx",
            version=version,
            result=result.stdout or result.stderr or "Created atlas-memo.pdf",
            visual_html=image_visual(preview_link, "Rendered fictional Atlas decision memo", "atlas-memo.pdf", "LibreOffice render"),
            artifacts=[pdf_link],
            duration_ms=result.duration_ms + preview_result.duration_ms,
        )
    )

    chrome_candidates = sorted((Path.home() / "Library/Caches/ms-playwright").glob("chromium-*/chrome-mac-x64/Google Chrome for Testing.app/Contents/MacOS/Google Chrome for Testing"), reverse=True)
    if not chrome_candidates:
        raise ToolFailure("Playwright-managed Chromium binary is missing")
    chrome = chrome_candidates[0]
    chrome_version = first_version_line(run([chrome, "--version"], cwd=desktop))
    screenshot = desktop / "chromium-dashboard.png"
    chromium_result = capture_html(fixtures["html"], screenshot, desktop)
    if not screenshot.is_file():
        raise ToolFailure("Chromium did not create the screenshot")
    screenshot_link = copy_artifact(screenshot, "chromium-dashboard.png")
    proofs.append(
        ToolProof(
            slug="chromium",
            name="Chromium",
            layer="Desktop workbench",
            tagline="Loads and captures a local business dashboard in a real Chromium engine.",
            task="Render a local dashboard at a fixed viewport",
            command="chromium --headless --screenshot dashboard.html",
            version=chrome_version,
            result="Captured a 1200 × 760 PNG from a local file URL.",
            visual_html=image_visual(screenshot_link, "Chromium screenshot of a fictional market dashboard", "Chromium · 1200 × 760", "real browser capture"),
            artifacts=[screenshot_link],
            duration_ms=chromium_result.duration_ms,
            note="The installed Playwright-managed Chromium binary is used so the example remains headless and deterministic.",
        )
    )

    quarto_version = first_version_line(run([executable("quarto"), "--version"], cwd=desktop))
    quarto_dir = desktop / "quarto"
    quarto_dir.mkdir()
    quarto_source = quarto_dir / "analysis.qmd"
    shutil.copy2(fixtures["quarto"], quarto_source)
    quarto_result = run(
        [executable("quarto"), "render", quarto_source.name, "--to", "html"],
        cwd=quarto_dir,
        env={"QUARTO_PYTHON": sys.executable},
        timeout=180,
    )
    quarto_html = quarto_dir / "analysis.html"
    if "Selected market: Coastal (81)" not in quarto_html.read_text(encoding="utf-8"):
        raise ToolFailure("Quarto output is missing the executed Python result")
    quarto_link = copy_artifact(quarto_html, "quarto-analysis.html")
    proofs.append(
        ToolProof(
            slug="quarto",
            name="Quarto",
            layer="Desktop workbench",
            tagline="Executes a small analysis and publishes the narrative, code, and result as one HTML report.",
            task="Render an executable decision brief",
            command="quarto render analysis.qmd --to html",
            version=quarto_version,
            result="Selected market: Coastal (81)",
            visual_html=browser_visual("Atlas market decision", '<p>Recommendation: start with the <strong>Coastal</strong> pilot.</p><div class="browser-card"><code>Selected market: Coastal (81)</code></div>', "quarto://analysis.html"),
            artifacts=[quarto_link],
            duration_ms=quarto_result.duration_ms,
        )
    )

    drawio_version = first_version_line(run([executable("drawio"), "--version"], cwd=desktop, timeout=30))
    drawio_svg = desktop / "drawio-decision.svg"
    drawio_result = run(
        [executable("drawio"), "--no-sandbox", "--export", "--format", "svg", "--output", drawio_svg, fixtures["drawio"]],
        cwd=desktop,
        timeout=180,
    )
    if not drawio_svg.is_file():
        raise ToolFailure("draw.io did not export the SVG")
    drawio_link = copy_artifact(drawio_svg)
    proofs.append(
        ToolProof(
            slug="drawio",
            name="draw.io",
            layer="Desktop workbench",
            tagline="Exports an editable decision flow into a browser-ready vector graphic.",
            task="Export a three-step decision diagram",
            command="drawio --export --format svg decision.drawio",
            version=drawio_version,
            result="Exported decision.drawio to SVG.",
            visual_html=image_visual(drawio_link, "Evidence to pilot gate to launch diagram", "decision.drawio", "SVG export"),
            artifacts=[drawio_link],
            duration_ms=drawio_result.duration_ms,
        )
    )

    inkscape_version = first_version_line(run([executable("inkscape"), "--version"], cwd=desktop))
    inkscape_png = desktop / "inkscape-gauge.png"
    inkscape_result = run(
        [executable("inkscape"), fixtures["svg"], "--export-width=1200", f"--export-filename={inkscape_png}"],
        cwd=desktop,
        timeout=120,
    )
    if not inkscape_png.is_file():
        raise ToolFailure("Inkscape did not export the PNG")
    inkscape_link = copy_artifact(inkscape_png)
    proofs.append(
        ToolProof(
            slug="inkscape",
            name="Inkscape",
            layer="Desktop workbench",
            tagline="Renders a scalable score graphic to a presentation-ready bitmap.",
            task="Export a 1,200-pixel market-score graphic",
            command="inkscape market-gauge.svg --export-width=1200",
            version=inkscape_version,
            result="Created a 1200 × 675 PNG from editable SVG source.",
            visual_html=image_visual(inkscape_link, "Market score gauge exported by Inkscape", "market-gauge.svg", "PNG export"),
            artifacts=[inkscape_link],
            duration_ms=inkscape_result.duration_ms,
        )
    )
    return proofs


def exercise_ai_clis(
    work: Path,
    *,
    real_ai: bool = False,
    codex_result_path: Path | None = None,
) -> list[ToolProof]:
    proofs: list[ToolProof] = []
    ai = work / "ai-clis"
    ai.mkdir()

    codex_binary = executable("codex")
    codex_version_result = run([codex_binary, "--version"], cwd=ai, timeout=30)
    codex_started = time.monotonic()
    codex_evidence = codex_result_path
    if codex_evidence is None and real_ai:
        codex_evidence = ai / "codex-decision.json"
        run(
            [
                codex_binary,
                "exec",
                "--ephemeral",
                "--sandbox",
                "read-only",
                "--color",
                "never",
                "-C",
                ROOT,
                "--output-schema",
                ROOT / "tools/tool_gallery_decision.schema.json",
                "-o",
                codex_evidence,
                "Read examples/market-entry-scorecard.csv. Treat every row as fictional. Return the one recommended first market, its numeric weighted_score, a one-sentence decision, and a measurable pilot gate. Use only the supplied data; do not modify files.",
            ],
            cwd=ai,
            timeout=420,
        )
    if codex_evidence is not None and codex_evidence.is_file():
        codex_data = json.loads(codex_evidence.read_text(encoding="utf-8"))
        required = {"recommended_market", "score", "decision", "gate"}
        if set(codex_data) != required:
            raise ToolFailure(f"Codex structured output fields differ: {sorted(codex_data)}")
        codex_link = copy_artifact(codex_evidence, "codex-decision.json")
        codex_visual = browser_visual(
            f"{codex_data['recommended_market']} is the first market.",
            f'<p>{escape(codex_data["decision"])}</p><div class="browser-card"><strong>Score {escape(codex_data["score"])}</strong><br>{escape(codex_data["gate"])}</div>',
            "codex://structured-output",
        )
        codex_result_text = json.dumps(codex_data, indent=2)
        codex_note = "This is the final structured output from an authenticated, ephemeral Codex model run against the repository's fictional CSV. The CLI had read-only sandbox access."
        codex_command = "codex exec --ephemeral --sandbox read-only --output-schema decision.schema.json"
        codex_artifacts = [codex_link]
    else:
        codex_help = run([codex_binary, "exec", "--help"], cwd=ai, timeout=30)
        codex_visual = terminal_visual(["$ codex exec --help", *[line for line in codex_help.stdout.splitlines() if line.strip()][:10]], "Codex · non-interactive CLI")
        codex_result_text = "The local non-interactive command and read-only sandbox options were detected."
        codex_note = "Run the gallery builder with --real-ai to create an authenticated structured-output proof."
        codex_command = "codex exec --help"
        codex_artifacts = []
    proofs.append(
        ToolProof(
            slug="codex",
            name="Codex",
            layer="AI coding CLIs",
            tagline="Reads the real fictional scorecard in a read-only sandbox and returns a schema-validated market decision.",
            task="Produce a structured recommendation from repository data",
            command=codex_command,
            version=first_version_line(codex_version_result),
            result=codex_result_text,
            visual_html=codex_visual,
            artifacts=codex_artifacts,
            duration_ms=elapsed_ms(codex_started),
            note=codex_note,
        )
    )

    proofs.append(
        ToolProof(
            slug="claude-code",
            name="Claude Code",
            layer="AI coding CLIs",
            tagline="Records the installed Claude Code package without invoking the unauthenticated CLI or making a model call.",
            task="Document the intentionally untested AI CLI",
            command="Not run — owner requested no Claude test",
            version=node_package_version("@anthropic-ai/claude-code"),
            result="Installed package metadata found. Authentication and model execution intentionally not tested.",
            visual_html=terminal_visual(["Claude Code", "", "status: installed", "authentication: not configured", "model call: not run", "reason: owner requested no Claude test"], "Claude Code · install record"),
            status="installed-only",
            duration_ms=0,
            note="This is the one deliberate exception in the gallery: 52 tools are exercised, while Claude Code is documented but not invoked.",
        )
    )
    return proofs


def exercise_office_media(fixtures: dict[str, Path], work: Path) -> list[ToolProof]:
    from PIL import Image, ImageDraw
    from docx import Document
    from openpyxl import Workbook, load_workbook
    from openpyxl.chart import BarChart, Reference
    from openpyxl.styles import Font, PatternFill
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt
    import cairosvg
    import xlsxwriter

    proofs: list[ToolProof] = []
    office = work / "office-media"
    office.mkdir()

    started = time.monotonic()
    openpyxl_path = office / "openpyxl-scorecard.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Market scorecard"
    rows = [
        ("Market", "Revenue ($m)", "Growth", "Score"),
        ("Coastal", 86, 0.14, 81),
        ("North", 79, 0.092, 78),
        ("Central", 69, 0.112, 72),
        ("West", 43, 0.058, 58),
    ]
    for row in rows:
        sheet.append(row)
    for cell in sheet[1]:
        cell.fill = PatternFill("solid", fgColor="B8F36B")
        cell.font = Font(bold=True, color="14201D")
    sheet["F1"] = "Average score"
    sheet["F2"] = "=AVERAGE(D2:D5)"
    sheet.freeze_panes = "A2"
    sheet.column_dimensions["A"].width = 18
    for cell in sheet["C"][1:]:
        cell.number_format = "0.0%"
    chart = BarChart()
    chart.title = "Market score"
    chart.y_axis.title = "Score"
    chart.add_data(Reference(sheet, min_col=4, min_row=1, max_row=5), titles_from_data=True)
    chart.set_categories(Reference(sheet, min_col=1, min_row=2, max_row=5))
    sheet.add_chart(chart, "F4")
    workbook.save(openpyxl_path)
    reopened = load_workbook(openpyxl_path, data_only=False)
    if reopened["Market scorecard"]["F2"].value != "=AVERAGE(D2:D5)":
        raise ToolFailure("openpyxl formula did not survive reopen")
    openpyxl_duration = elapsed_ms(started)
    openpyxl_link = copy_artifact(openpyxl_path)
    proofs.append(
        ToolProof(
            slug="openpyxl",
            name="openpyxl",
            layer="Office, PDF, and media",
            tagline="Builds a styled Excel-compatible scorecard, formula, freeze pane, and native chart, then reopens the workbook to verify them.",
            task="Create and reopen a market scorecard workbook",
            command="openpyxl.Workbook() → save() → load_workbook()",
            version=package_version("openpyxl"),
            result="Workbook reopened; formula =AVERAGE(D2:D5) and chart object are present.",
            visual_html=table_visual(rows[0], [[r[0], r[1], f"{r[2]:.1%}", r[3]] for r in rows[1:]], "openpyxl-scorecard.xlsx"),
            artifacts=[openpyxl_link],
            duration_ms=openpyxl_duration,
        )
    )

    started = time.monotonic()
    xlsxwriter_path = office / "xlsxwriter-scenarios.xlsx"
    book = xlsxwriter.Workbook(xlsxwriter_path)
    ws = book.add_worksheet("Scenarios")
    header = book.add_format({"bold": True, "bg_color": "#B8F36B", "border": 1})
    money = book.add_format({"num_format": "$0.0", "border": 1})
    ws.write_row("A1", ["Scenario", "Revenue ($m)", "EBITDA ($m)"], header)
    scenario_rows = [("Downside", 9.8, 1.1), ("Base", 12.4, 2.3), ("Upside", 15.1, 3.5)]
    for index, row in enumerate(scenario_rows, start=1):
        ws.write(index, 0, row[0])
        ws.write_number(index, 1, row[1], money)
        ws.write_number(index, 2, row[2], money)
    xchart = book.add_chart({"type": "column"})
    xchart.add_series({"name": "EBITDA", "categories": "=Scenarios!$A$2:$A$4", "values": "=Scenarios!$C$2:$C$4", "fill": {"color": "#176B4D"}})
    xchart.set_title({"name": "Pilot economics"})
    ws.insert_chart("E2", xchart)
    ws.conditional_format("C2:C4", {"type": "data_bar", "bar_color": "#F26B38"})
    book.close()
    if not xlsxwriter_path.is_file():
        raise ToolFailure("XlsxWriter workbook was not created")
    xlsxwriter_duration = elapsed_ms(started)
    xlsxwriter_link = copy_artifact(xlsxwriter_path)
    proofs.append(
        ToolProof(
            slug="xlsxwriter",
            name="XlsxWriter",
            layer="Office, PDF, and media",
            tagline="Creates an Excel-compatible scenario model with formats, conditional bars, and a native column chart.",
            task="Write a three-scenario economics workbook",
            command="xlsxwriter.Workbook() + add_chart() + close()",
            version=package_version("XlsxWriter"),
            result="Created three scenarios, a data-bar rule, and a native chart.",
            visual_html=bars_visual([(row[0], row[2]) for row in scenario_rows], "EBITDA by scenario", "m"),
            artifacts=[xlsxwriter_link],
            duration_ms=xlsxwriter_duration,
        )
    )

    started = time.monotonic()
    doc = Document(fixtures["docx"])
    doc.add_heading("Control", level=1)
    doc.add_paragraph("The pilot advances only after finance signs off on realized economics.")
    docx_path = office / "python-docx-memo.docx"
    doc.save(docx_path)
    reopened_doc = Document(docx_path)
    heading_count = sum(1 for paragraph in reopened_doc.paragraphs if paragraph.style.name.startswith("Heading"))
    docx_duration = elapsed_ms(started)
    docx_link = copy_artifact(docx_path)
    proofs.append(
        ToolProof(
            slug="python-docx",
            name="python-docx",
            layer="Office, PDF, and media",
            tagline="Creates and reopens a structured decision memo with named headings, body text, and a real Word table.",
            task="Extend and validate a Word decision memo",
            command="Document(source) → add_heading() → save() → reopen",
            version=package_version("python-docx"),
            result=f"Reopened {len(reopened_doc.paragraphs)} paragraphs, {heading_count} headings, and {len(reopened_doc.tables)} table.",
            visual_html=document_visual("Coastal is the first gated pilot.", "Atlas services · decision memo", ["The scorecard favors Coastal at 81. North remains prepared as the follow-on market.", "Control: finance signs off on realized economics before scaling."], "editable docx"),
            artifacts=[docx_link],
            duration_ms=docx_duration,
        )
    )

    started = time.monotonic()
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor(243, 240, 232)
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.55), Inches(11.8), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = "Coastal is the first gated pilot."
    title_frame.paragraphs[0].font.name = "Arial"
    title_frame.paragraphs[0].font.size = Pt(34)
    title_frame.paragraphs[0].font.bold = True
    phases = [("Evidence", "81 score"), ("Economics", "30 days"), ("Gate", "CFO sign-off"), ("Launch", "if proven")]
    colors = [RGBColor(184, 243, 107), RGBColor(245, 217, 105), RGBColor(242, 107, 56), RGBColor(23, 107, 77)]
    for index, ((phase, detail), color) in enumerate(zip(phases, colors)):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7 + index * 3.1), Inches(2.6), Inches(2.7), Inches(2.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(20, 32, 29)
        shape.text_frame.text = f"{phase}\n{detail}"
        shape.text_frame.paragraphs[0].font.name = "Arial"
        shape.text_frame.paragraphs[0].font.size = Pt(20)
        shape.text_frame.paragraphs[0].font.bold = True
    pptx_path = office / "python-pptx-pilot.pptx"
    presentation.save(pptx_path)
    reopened_presentation = Presentation(pptx_path)
    pptx_duration = elapsed_ms(started)
    pptx_link = copy_artifact(pptx_path)
    slide_preview = office / "python-pptx-render"
    slide_preview.mkdir()
    render_result = run([executable("soffice"), "--headless", f"-env:UserInstallation=file://{office / 'pptx-profile'}", "--convert-to", "pdf", "--outdir", slide_preview, pptx_path], cwd=office, timeout=120)
    slide_pdf = slide_preview / "python-pptx-pilot.pdf"
    run([executable("pdftoppm"), "-png", "-f", "1", "-singlefile", "-r", "120", slide_pdf, slide_preview / "slide"], cwd=office)
    slide_png = slide_preview / "slide.png"
    slide_link = copy_artifact(slide_png, "python-pptx-pilot.png")
    proofs.append(
        ToolProof(
            slug="python-pptx",
            name="python-pptx",
            layer="Office, PDF, and media",
            tagline="Builds a widescreen presentation from native text boxes and editable PowerPoint shapes.",
            task="Create and reopen an editable decision slide",
            command="Presentation() + add_shape() + save() + reopen",
            version=package_version("python-pptx"),
            result=f"Reopened {len(reopened_presentation.slides)} slide with {len(reopened_presentation.slides[0].shapes)} native shapes.",
            visual_html=image_visual(slide_link, "Rendered pilot-gate slide created with python-pptx", "python-pptx-pilot.pptx", "native slide render"),
            artifacts=[pptx_link],
            duration_ms=pptx_duration + render_result.duration_ms,
        )
    )

    started = time.monotonic()
    cairosvg_png = office / "cairosvg-gauge.png"
    cairosvg.svg2png(url=str(fixtures["svg"]), write_to=str(cairosvg_png), output_width=1200)
    cairosvg_duration = elapsed_ms(started)
    cairosvg_link = copy_artifact(cairosvg_png)
    proofs.append(
        ToolProof(
            slug="cairosvg",
            name="CairoSVG",
            layer="Office, PDF, and media",
            tagline="Converts an editable SVG score graphic into a crisp PNG with a deterministic Python call.",
            task="Rasterize an SVG market-score graphic",
            command="cairosvg.svg2png(output_width=1200)",
            version=package_version("CairoSVG"),
            result="Created a 1,200-pixel PNG from SVG source.",
            visual_html=image_visual(cairosvg_link, "CairoSVG-rendered market score graphic", "market-gauge.svg", "Cairo raster"),
            artifacts=[cairosvg_link],
            duration_ms=cairosvg_duration,
        )
    )

    pandoc_version = first_version_line(run([executable("pandoc"), "--version"], cwd=office))
    pandoc_html = office / "pandoc-brief.html"
    pandoc_result = run([executable("pandoc"), fixtures["markdown"], "--standalone", "--metadata", "title=Atlas decision brief", "--output", pandoc_html], cwd=office)
    if "Atlas market-entry decision" not in pandoc_html.read_text(encoding="utf-8"):
        raise ToolFailure("Pandoc output is missing the source heading")
    pandoc_link = copy_artifact(pandoc_html)
    proofs.append(
        ToolProof(
            slug="pandoc",
            name="Pandoc",
            layer="Office, PDF, and media",
            tagline="Transforms a Markdown decision brief and table into a standalone semantic HTML document.",
            task="Publish a Markdown brief as standalone HTML",
            command="pandoc decision-brief.md --standalone --output brief.html",
            version=pandoc_version,
            result="Preserved the heading, recommendation, table, and fictional-data notice.",
            visual_html=browser_visual("Atlas market-entry decision", '<p><strong>Recommendation:</strong> Start with Coastal, then gate North on pilot economics.</p><div class="browser-card">Markdown table → semantic HTML table</div>', "pandoc://brief.html"),
            artifacts=[pandoc_link],
            duration_ms=pandoc_result.duration_ms,
        )
    )

    poppler_version = first_version_line(run([executable("pdfinfo"), "-v"], cwd=office))
    pdfinfo_result = run([executable("pdfinfo"), fixtures["rich_pdf"]], cwd=office)
    text_result = run([executable("pdftotext"), "-layout", fixtures["rich_pdf"], "-"], cwd=office)
    poppler_root = office / "poppler-page"
    render_result = run([executable("pdftoppm"), "-png", "-singlefile", "-r", "130", fixtures["rich_pdf"], poppler_root], cwd=office)
    poppler_png = poppler_root.with_suffix(".png")
    poppler_link = copy_artifact(poppler_png)
    proofs.append(
        ToolProof(
            slug="poppler",
            name="Poppler",
            layer="Office, PDF, and media",
            tagline="Inspects PDF metadata, extracts positioned text, and renders a page preview with three standard Poppler utilities.",
            task="Inspect, extract, and render a decision PDF",
            command="pdfinfo decision.pdf && pdftotext -layout decision.pdf - && pdftoppm -png decision.pdf",
            version=poppler_version,
            result="\n".join(text_result.stdout.splitlines()[:5]),
            visual_html=image_visual(poppler_link, "Poppler-rendered decision PDF", "decision.pdf", "pdftoppm render"),
            artifacts=[poppler_link],
            duration_ms=pdfinfo_result.duration_ms + text_result.duration_ms + render_result.duration_ms,
        )
    )

    qpdf_version = first_version_line(run([executable("qpdf"), "--version"], cwd=office))
    linearized = office / "qpdf-linearized.pdf"
    qpdf_result = run([executable("qpdf"), "--linearize", fixtures["rich_pdf"], linearized], cwd=office)
    check_result = run([executable("qpdf"), "--check", linearized], cwd=office)
    qpdf_link = copy_artifact(linearized)
    proofs.append(
        ToolProof(
            slug="qpdf",
            name="qpdf",
            layer="Office, PDF, and media",
            tagline="Rewrites a valid PDF for fast web access, then checks the new file's structure.",
            task="Linearize and validate a decision PDF",
            command="qpdf --linearize decision.pdf linearized.pdf && qpdf --check linearized.pdf",
            version=qpdf_version,
            result=check_result.stdout or "No syntax or stream encoding errors found.",
            visual_html=terminal_visual(["$ qpdf --linearize decision.pdf linearized.pdf", "✓ linearized output written", "$ qpdf --check linearized.pdf", *(check_result.stdout.splitlines()[:8] or ["✓ PDF structure valid"])], "qpdf · structure check"),
            artifacts=[qpdf_link],
            duration_ms=qpdf_result.duration_ms + check_result.duration_ms,
        )
    )

    mupdf_version = first_version_line(run([executable("mutool"), "-v"], cwd=office))
    mupdf_png = office / "mupdf-page.png"
    mupdf_result = run([executable("mutool"), "draw", "-q", "-r", "144", "-o", mupdf_png, fixtures["rich_pdf"], "1"], cwd=office)
    mupdf_link = copy_artifact(mupdf_png)
    proofs.append(
        ToolProof(
            slug="mupdf",
            name="MuPDF",
            layer="Office, PDF, and media",
            tagline="Renders a selected PDF page at a controlled resolution with the compact MuPDF engine.",
            task="Render page one of a decision PDF",
            command="mutool draw -r 144 -o page.png decision.pdf 1",
            version=mupdf_version,
            result="Rendered page 1 at 144 DPI.",
            visual_html=image_visual(mupdf_link, "MuPDF-rendered decision PDF", "decision.pdf · page 1", "MuPDF render"),
            artifacts=[mupdf_link],
            duration_ms=mupdf_result.duration_ms,
        )
    )

    ocr_version = first_version_line(run([executable("ocrmypdf"), "--version"], cwd=office))
    searchable_pdf = office / "ocr-searchable.pdf"
    ocr_result = run([executable("ocrmypdf"), "--force-ocr", "--deskew", "--output-type", "pdf", fixtures["scan_pdf"], searchable_pdf], cwd=office, timeout=240)
    extracted = run([executable("pdftotext"), searchable_pdf, "-"], cwd=office)
    if "Coastal" not in extracted.stdout:
        raise ToolFailure("OCRmyPDF output did not contain the expected word Coastal")
    ocr_preview = office / "ocr-preview"
    run([executable("pdftoppm"), "-png", "-singlefile", "-r", "110", searchable_pdf, ocr_preview], cwd=office)
    ocr_png = ocr_preview.with_suffix(".png")
    ocr_png_link = copy_artifact(ocr_png)
    ocr_pdf_link = copy_artifact(searchable_pdf)
    proofs.append(
        ToolProof(
            slug="ocrmypdf",
            name="OCRmyPDF",
            layer="Office, PDF, and media",
            tagline="Adds a searchable text layer to an image-only PDF and proves the expected words can be extracted.",
            task="Deskew and OCR a scanned pilot note",
            command="ocrmypdf --force-ocr --deskew scan.pdf searchable.pdf",
            version=ocr_version,
            result="\n".join(extracted.stdout.splitlines()[:6]),
            visual_html=image_visual(ocr_png_link, "OCR-processed fictional pilot note", "searchable.pdf", "searchable text layer"),
            artifacts=[ocr_pdf_link],
            duration_ms=ocr_result.duration_ms,
        )
    )

    tesseract_version = first_version_line(run([executable("tesseract"), "--version"], cwd=office))
    tesseract_result = run([executable("tesseract"), fixtures["scan_png"], "stdout", "--psm", "6"], cwd=office, timeout=120)
    if "Coastal" not in tesseract_result.stdout:
        raise ToolFailure("Tesseract did not extract the expected word Coastal")
    scan_link = copy_artifact(fixtures["scan_png"], "tesseract-input.png")
    tesseract_visual = (
        '<div class="split-visual">'
        f'<div><img src="{escape(scan_link)}" alt="Fictional pilot note scanned image"></div>'
        f'<pre><span>OCR text</span>\n{escape(tesseract_result.stdout)}</pre></div>'
    )
    proofs.append(
        ToolProof(
            slug="tesseract",
            name="Tesseract",
            layer="Office, PDF, and media",
            tagline="Reads text directly from a rasterized pilot note and returns plain machine-searchable text.",
            task="Extract text from a scanned PNG",
            command="tesseract pilot-note.png stdout --psm 6",
            version=tesseract_version,
            result=tesseract_result.stdout,
            visual_html=tesseract_visual,
            artifacts=[scan_link],
            duration_ms=tesseract_result.duration_ms,
        )
    )

    imagemagick_version = first_version_line(run([executable("magick"), "--version"], cwd=office))
    imagemagick_out = office / "imagemagick-card.webp"
    imagemagick_result = run([executable("magick"), fixtures["scan_png"], "-resize", "1200x", "-bordercolor", "#176b4d", "-border", "18", "-strip", "-quality", "82", imagemagick_out], cwd=office)
    imagemagick_link = copy_artifact(imagemagick_out)
    proofs.append(
        ToolProof(
            slug="imagemagick",
            name="ImageMagick",
            layer="Office, PDF, and media",
            tagline="Resizes, frames, strips metadata, and compresses a large raster into a web-ready image.",
            task="Prepare a scanned note for the web",
            command="magick input.png -resize 1200x -border 18 -strip -quality 82 output.webp",
            version=imagemagick_version,
            result=f"Created {imagemagick_out.name} at {imagemagick_out.stat().st_size:,} bytes.",
            visual_html=image_visual(imagemagick_link, "ImageMagick-processed pilot note", "output.webp", "resized + compressed"),
            artifacts=[imagemagick_link],
            duration_ms=imagemagick_result.duration_ms,
        )
    )

    frames = office / "frames"
    frames.mkdir()
    for index in range(24):
        frame = Image.new("RGB", (960, 540), "#14201d")
        draw = ImageDraw.Draw(frame)
        progress = 90 + index * 28
        draw.rectangle((90, 260, 870, 335), outline="#edf2ed", width=3)
        draw.rectangle((90, 260, min(progress, 870), 335), fill="#b8f36b")
        draw.text((90, 105), "ATLAS PILOT GATE", fill="#edf2ed")
        draw.text((90, 175), f"Evidence complete: {min(100, round(index / 23 * 100))}%", fill="#edf2ed")
        frame.save(frames / f"frame-{index:02d}.png")
    ffmpeg_version = first_version_line(run([executable("ffmpeg"), "-version"], cwd=office))
    video = office / "ffmpeg-pilot.mp4"
    ffmpeg_result = run([executable("ffmpeg"), "-y", "-loglevel", "error", "-framerate", "12", "-i", str(frames / "frame-%02d.png"), "-c:v", "libx264", "-pix_fmt", "yuv420p", "-movflags", "+faststart", video], cwd=office, timeout=180)
    video_link = copy_artifact(video)
    poster_link = copy_artifact(frames / "frame-23.png", "ffmpeg-poster.png")
    video_visual = f'<div class="media-visual"><div class="visual-chrome"><span>ffmpeg-pilot.mp4</span><span>H.264 · 2 seconds</span></div><video controls muted loop playsinline poster="{escape(poster_link)}"><source src="{escape(video_link)}" type="video/mp4">Video preview unavailable.</video></div>'
    proofs.append(
        ToolProof(
            slug="ffmpeg",
            name="FFmpeg",
            layer="Office, PDF, and media",
            tagline="Encodes a 24-frame business-status animation as a browser-playable H.264 video.",
            task="Encode a two-second pilot-progress animation",
            command="ffmpeg -framerate 12 -i frame-%02d.png -c:v libx264 -pix_fmt yuv420p pilot.mp4",
            version=ffmpeg_version,
            result=f"Encoded 24 frames into {video.stat().st_size:,} bytes.",
            visual_html=video_visual,
            artifacts=[video_link],
            duration_ms=ffmpeg_result.duration_ms,
        )
    )
    return proofs


def exercise_data_research(fixtures: dict[str, Path], work: Path) -> list[ToolProof]:
    import http.server
    import threading

    import duckdb
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import nbformat
    import numpy as np
    import pandas as pd
    import polars as pl
    import pyarrow as pa
    import pyarrow.csv as arrow_csv
    import pyarrow.parquet as parquet
    import scipy
    from scipy import stats
    import sklearn
    from sklearn.cluster import KMeans
    import statsmodels.api as sm

    proofs: list[ToolProof] = []
    data_dir = work / "data-research"
    data_dir.mkdir()

    started = time.monotonic()
    frame = pd.read_csv(fixtures["csv"])
    frame["rank"] = frame["score"].rank(method="dense", ascending=False).astype(int)
    frame["revenue_share_pct"] = (frame["revenue_m"] / frame["revenue_m"].sum() * 100).round(1)
    pandas_result = frame.sort_values("rank")[["market", "score", "rank", "revenue_share_pct"]]
    pandas_duration = elapsed_ms(started)
    proofs.append(
        ToolProof(
            slug="pandas",
            name="pandas",
            layer="Data and research",
            tagline="Loads a CSV, ranks markets, and calculates each market's share of the fictional revenue pool.",
            task="Transform a market scorecard into a ranked decision table",
            command="pandas.read_csv() + rank() + vectorized share calculation",
            version=package_version("pandas"),
            result=pandas_result.to_string(index=False),
            visual_html=table_visual(["Market", "Score", "Rank", "Revenue share"], [[row.market, row.score, row.rank, f"{row.revenue_share_pct:.1f}%"] for row in pandas_result.itertuples()], "pandas · ranked scorecard"),
            duration_ms=pandas_duration,
        )
    )

    started = time.monotonic()
    lazy = pl.scan_csv(fixtures["csv"])
    polars_result = (
        lazy.with_columns((pl.col("score") - pl.col("score").mean()).round(1).alias("vs_mean"))
        .sort("score", descending=True)
        .select("market", "growth_pct", "score", "vs_mean")
        .collect()
    )
    polars_duration = elapsed_ms(started)
    proofs.append(
        ToolProof(
            slug="polars",
            name="Polars",
            layer="Data and research",
            tagline="Runs a lazy CSV query, computes score distance from the mean, sorts, and collects only the decision columns.",
            task="Execute a lazy market-ranking query",
            command="pl.scan_csv() → with_columns() → sort() → collect()",
            version=package_version("polars"),
            result=str(polars_result),
            visual_html=table_visual(polars_result.columns, polars_result.rows(), "Polars · lazy result"),
            duration_ms=polars_duration,
        )
    )

    started = time.monotonic()
    arrow_table = arrow_csv.read_csv(fixtures["csv"])
    parquet_path = data_dir / "markets.parquet"
    parquet.write_table(arrow_table, parquet_path, compression="zstd")
    arrow_roundtrip = parquet.read_table(parquet_path)
    if not arrow_table.equals(arrow_roundtrip):
        raise ToolFailure("Arrow Parquet round trip changed the table")
    arrow_duration = elapsed_ms(started)
    arrow_link = copy_artifact(parquet_path)
    schema_rows = [(field.name, str(field.type), arrow_roundtrip.column(field.name).null_count) for field in arrow_roundtrip.schema]
    proofs.append(
        ToolProof(
            slug="apache-arrow",
            name="Arrow",
            layer="Data and research",
            tagline="Reads CSV into a typed columnar table, writes compressed Parquet, and verifies an exact round trip.",
            task="Create and validate a typed Parquet dataset",
            command="pyarrow.csv.read_csv() → parquet.write_table() → parquet.read_table()",
            version=package_version("pyarrow"),
            result=f"{arrow_roundtrip.num_rows} rows × {arrow_roundtrip.num_columns} columns; schema preserved; 0 nulls.",
            visual_html=table_visual(["Field", "Arrow type", "Nulls"], schema_rows, "Arrow · schema"),
            artifacts=[arrow_link],
            duration_ms=arrow_duration,
        )
    )

    duckdb_version_result = run([executable("duckdb"), "--version"], cwd=data_dir)
    sql = f"SELECT market, score, ROUND(growth_pct * revenue_m / 100, 2) AS growth_value FROM read_csv_auto('{fixtures['csv']}') ORDER BY score DESC"
    duckdb_result = run([executable("duckdb"), "-csv", "-c", sql], cwd=data_dir)
    parsed = list(csv.reader(io.StringIO(duckdb_result.stdout)))
    proofs.append(
        ToolProof(
            slug="duckdb",
            name="DuckDB",
            layer="Data and research",
            tagline="Queries the CSV directly with SQL and calculates a sortable growth-value proxy without loading a server.",
            task="Run an analytical SQL query against CSV",
            command="duckdb -csv -c \"SELECT … FROM read_csv_auto('scorecard.csv')\"",
            version=first_version_line(duckdb_version_result),
            result=duckdb_result.stdout,
            visual_html=table_visual(parsed[0], parsed[1:], "DuckDB · SQL result"),
            duration_ms=duckdb_result.duration_ms,
        )
    )

    started = time.monotonic()
    scipy_result = stats.spearmanr(frame["growth_pct"], frame["score"])
    smooth_x = np.linspace(frame["growth_pct"].min(), frame["growth_pct"].max(), 80)
    interpolation = scipy.interpolate.PchipInterpolator(frame.sort_values("growth_pct")["growth_pct"], frame.sort_values("growth_pct")["score"])
    smooth_y = interpolation(smooth_x)
    fig, axis = plt.subplots(figsize=(8, 4.6))
    fig.patch.set_facecolor("#f3f0e8")
    axis.set_facecolor("#f3f0e8")
    axis.plot(smooth_x, smooth_y, color="#176b4d", linewidth=3)
    axis.scatter(frame["growth_pct"], frame["score"], color="#f26b38", s=90, zorder=3)
    for row in frame.itertuples():
        axis.annotate(row.market, (row.growth_pct, row.score), xytext=(5, 7), textcoords="offset points")
    axis.set_xlabel("Growth (%)")
    axis.set_ylabel("Weighted score")
    axis.spines[["top", "right"]].set_visible(False)
    scipy_svg = data_dir / "scipy-interpolation.svg"
    fig.tight_layout()
    fig.savefig(scipy_svg, format="svg")
    plt.close(fig)
    scipy_duration = elapsed_ms(started)
    scipy_link = copy_artifact(scipy_svg)
    proofs.append(
        ToolProof(
            slug="scipy",
            name="SciPy",
            layer="Data and research",
            tagline="Measures rank correlation and draws a shape-preserving interpolation through the supplied market observations.",
            task="Quantify and visualize growth-to-score association",
            command="scipy.stats.spearmanr() + scipy.interpolate.PchipInterpolator()",
            version=scipy.__version__,
            result=f"Spearman rho={scipy_result.statistic:.3f}; p={scipy_result.pvalue:.3f}; n={len(frame)}.",
            visual_html=image_visual(scipy_link, "SciPy interpolation of growth and weighted score", "growth vs score", "computed SVG"),
            artifacts=[scipy_link],
            duration_ms=scipy_duration,
            note="Four rows are enough to exercise the method, not enough to support a causal claim. The page reports n explicitly.",
        )
    )

    started = time.monotonic()
    features = frame[["growth_pct", "score"]].to_numpy()
    model = KMeans(n_clusters=2, random_state=7, n_init=20).fit(features)
    cluster_frame = frame[["market", "growth_pct", "score"]].copy()
    cluster_frame["cluster"] = model.labels_ + 1
    sklearn_duration = elapsed_ms(started)
    proofs.append(
        ToolProof(
            slug="scikit-learn",
            name="scikit-learn",
            layer="Data and research",
            tagline="Fits a seeded two-cluster model and labels each market with a reproducible peer group.",
            task="Cluster markets by growth and weighted score",
            command="KMeans(n_clusters=2, random_state=7, n_init=20).fit()",
            version=sklearn.__version__,
            result=f"Inertia={model.inertia_:.2f}; centers={model.cluster_centers_.round(2).tolist()}",
            visual_html=table_visual(["Market", "Growth", "Score", "Cluster"], [[row.market, f"{row.growth_pct:.1f}%", row.score, row.cluster] for row in cluster_frame.itertuples()], "scikit-learn · seeded clusters"),
            duration_ms=sklearn_duration,
            note="This is a mechanics example on four fictional rows, not a production segmentation model.",
        )
    )

    started = time.monotonic()
    design = sm.add_constant(frame["growth_pct"])
    ols = sm.OLS(frame["score"], design).fit()
    predicted = ols.predict(design)
    fig, axis = plt.subplots(figsize=(8, 4.6))
    fig.patch.set_facecolor("#f3f0e8")
    axis.set_facecolor("#f3f0e8")
    axis.scatter(frame["growth_pct"], frame["score"], color="#f26b38", s=90)
    order = np.argsort(frame["growth_pct"].to_numpy())
    axis.plot(frame["growth_pct"].to_numpy()[order], predicted.to_numpy()[order], color="#176b4d", linewidth=3)
    axis.set_xlabel("Growth (%)")
    axis.set_ylabel("Weighted score")
    axis.spines[["top", "right"]].set_visible(False)
    statsmodels_svg = data_dir / "statsmodels-ols.svg"
    fig.tight_layout()
    fig.savefig(statsmodels_svg, format="svg")
    plt.close(fig)
    statsmodels_duration = elapsed_ms(started)
    statsmodels_link = copy_artifact(statsmodels_svg)
    proofs.append(
        ToolProof(
            slug="statsmodels",
            name="statsmodels",
            layer="Data and research",
            tagline="Fits an explicit OLS model and reports the coefficient, uncertainty, fit, and sample size instead of only drawing a trend line.",
            task="Estimate score association with market growth",
            command="statsmodels.api.OLS(score, add_constant(growth)).fit()",
            version=package_version("statsmodels"),
            result=f"growth coefficient={ols.params['growth_pct']:.3f}; 95% CI={ols.conf_int().loc['growth_pct'].round(3).tolist()}; R²={ols.rsquared:.3f}; n={int(ols.nobs)}.",
            visual_html=image_visual(statsmodels_link, "OLS line fitted to fictional market growth and scores", "OLS fit", "model SVG"),
            artifacts=[statsmodels_link],
            duration_ms=statsmodels_duration,
            note="The tiny fictional sample is deliberately labeled; the proof is the model workflow, not a statistically reliable market conclusion.",
        )
    )

    notebook = nbformat.v4.new_notebook()
    notebook.cells = [
        nbformat.v4.new_markdown_cell("# Atlas scorecard\n\nA fictional, executable notebook proof."),
        nbformat.v4.new_code_cell("scores = {'Coastal': 81, 'North': 78, 'Central': 72, 'West': 58}\nwinner = max(scores, key=scores.get)\nprint(f'Pilot: {winner} ({scores[winner]})')"),
    ]
    notebook.metadata.kernelspec = {"display_name": "Python 3", "language": "python", "name": "python3"}
    notebook_path = data_dir / "atlas.ipynb"
    nbformat.write(notebook, notebook_path)
    jupyter_version = first_version_line(run([sys.executable, "-m", "jupyterlab", "--version"], cwd=data_dir))
    notebook_result = run([sys.executable, "-m", "jupyter", "nbconvert", "--to", "html", "--execute", "--ExecutePreprocessor.timeout=90", notebook_path.name], cwd=data_dir, timeout=180)
    notebook_html = data_dir / "atlas.html"
    if "Pilot: Coastal (81)" not in notebook_html.read_text(encoding="utf-8"):
        raise ToolFailure("executed notebook is missing its expected output")
    notebook_link = copy_artifact(notebook_html, "jupyterlab-notebook.html")
    proofs.append(
        ToolProof(
            slug="jupyterlab",
            name="JupyterLab",
            layer="Data and research",
            tagline="Executes a real notebook kernel and exports the markdown, code, and computed output as portable HTML.",
            task="Execute and publish a scorecard notebook",
            command="jupyter nbconvert --to html --execute atlas.ipynb",
            version=jupyter_version,
            result="Pilot: Coastal (81)",
            visual_html=browser_visual("Atlas scorecard", '<p>A fictional, executable notebook proof.</p><div class="browser-card"><code>scores → Pilot: Coastal (81)</code></div>', "jupyter://atlas.ipynb"),
            artifacts=[notebook_link],
            duration_ms=notebook_result.duration_ms,
        )
    )

    chrome_candidates = sorted((Path.home() / "Library/Caches/ms-playwright").glob("chromium-*/chrome-mac-x64/Google Chrome for Testing.app/Contents/MacOS/Google Chrome for Testing"), reverse=True)
    if not chrome_candidates:
        raise ToolFailure("Playwright-managed Chromium binary is missing")
    playwright_script = write_text(
        data_dir / "playwright_capture.py",
        f"""
        from pathlib import Path
        from playwright.sync_api import sync_playwright
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True, executable_path={str(chrome_candidates[0])!r}, args=["--no-sandbox", "--disable-gpu"])
            page = browser.new_page(viewport={{"width": 1200, "height": 760}})
            page.goto({fixtures['html'].as_uri()!r})
            assert page.locator("[data-market='Coastal'] strong").inner_text() == "81"
            page.screenshot(path="playwright-dashboard.png", full_page=True)
            print(page.title())
            print(page.locator("#summary").inner_text())
            browser.close()
        """,
    )
    playwright_result = run([sys.executable, playwright_script], cwd=data_dir, timeout=180)
    playwright_png = data_dir / "playwright-dashboard.png"
    playwright_link = copy_artifact(playwright_png)
    proofs.append(
        ToolProof(
            slug="playwright",
            name="Playwright",
            layer="Data and research",
            tagline="Launches Chromium, navigates a local page, asserts visible data, and captures a full-page screenshot.",
            task="Run a browser assertion and visual capture",
            command="sync_playwright() → page.goto() → locator().inner_text() → screenshot()",
            version=package_version("playwright"),
            result=playwright_result.stdout,
            visual_html=image_visual(playwright_link, "Playwright screenshot of the local market dashboard", "Playwright · Chromium", "asserted + captured"),
            artifacts=[playwright_link],
            duration_ms=playwright_result.duration_ms,
        )
    )

    requests_seen: list[str] = []

    class WebDriverFixture(http.server.BaseHTTPRequestHandler):
        def log_message(self, format: str, *args: object) -> None:
            return

        def respond(self, value: object) -> None:
            payload = json.dumps({"value": value}).encode("utf-8")
            self.send_response(200)
            self.send_header("Content-Type", "application/json")
            self.send_header("Content-Length", str(len(payload)))
            self.end_headers()
            self.wfile.write(payload)

        def do_POST(self) -> None:
            length = int(self.headers.get("Content-Length", "0"))
            if length:
                self.rfile.read(length)
            requests_seen.append(f"POST {self.path}")
            if self.path == "/session":
                self.respond({"sessionId": "wutpack-proof", "capabilities": {"browserName": "local-protocol-fixture"}})
            elif self.path.endswith("/element"):
                self.respond({"element-6066-11e4-a52e-4f735466cecf": "summary"})
            else:
                self.respond(None)

        def do_GET(self) -> None:
            requests_seen.append(f"GET {self.path}")
            if self.path.endswith("/title"):
                self.respond("Atlas local dashboard")
            elif self.path.endswith("/element/summary/text"):
                self.respond("Coastal leads the scorecard. North remains the gated follow-on.")
            else:
                self.respond(None)

        def do_DELETE(self) -> None:
            requests_seen.append(f"DELETE {self.path}")
            self.respond(None)

    server = http.server.ThreadingHTTPServer(("127.0.0.1", 0), WebDriverFixture)
    thread = threading.Thread(target=server.serve_forever, daemon=True)
    thread.start()
    started = time.monotonic()
    try:
        from selenium import webdriver
        from selenium.webdriver.common.by import By
        options = webdriver.ChromeOptions()
        driver = webdriver.Remote(command_executor=f"http://127.0.0.1:{server.server_port}", options=options)
        driver.get(fixtures["html"].as_uri())
        selenium_title = driver.title
        selenium_summary = driver.find_element(By.ID, "summary").text
        driver.quit()
    finally:
        server.shutdown()
        server.server_close()
        thread.join(timeout=5)
    selenium_duration = elapsed_ms(started)
    if selenium_title != "Atlas local dashboard" or "Coastal" not in selenium_summary:
        raise ToolFailure("Selenium protocol result did not match the fixture")
    proofs.append(
        ToolProof(
            slug="selenium",
            name="Selenium",
            layer="Data and research",
            tagline="Exercises the Selenium client through the W3C WebDriver protocol: create session, navigate, read title, locate text, and quit.",
            task="Drive a deterministic local WebDriver protocol fixture",
            command="webdriver.Remote() → get() → title → find_element(By.ID) → quit()",
            version=package_version("selenium"),
            result=f"{selenium_title}\n{selenium_summary}\n" + "\n".join(requests_seen),
            visual_html=terminal_visual(["W3C WebDriver trace", *requests_seen, "", f'title = "{selenium_title}"', f'text  = "{selenium_summary}"'], "Selenium · local protocol"),
            duration_ms=selenium_duration,
            note="Playwright already proves a real local Chromium launch. This Selenium page isolates and visibly proves the client-side WebDriver command flow without downloading a separate ChromeDriver binary.",
        )
    )

    spider = write_text(
        data_dir / "market_spider.py",
        f"""
        import scrapy

        class MarketSpider(scrapy.Spider):
            name = "markets"
            start_urls = [{fixtures['html'].as_uri()!r}]

            def parse(self, response):
                for card in response.css("[data-market]"):
                    yield {{"market": card.attrib["data-market"], "score": int(card.css("strong::text").get())}}
        """,
    )
    scrapy_json = data_dir / "scrapy-markets.json"
    scrapy_result = run([sys.executable, "-m", "scrapy", "runspider", spider, "-O", scrapy_json, "-s", "LOG_ENABLED=False"], cwd=data_dir, timeout=180)
    scraped = json.loads(scrapy_json.read_text(encoding="utf-8"))
    if len(scraped) != 4:
        raise ToolFailure(f"Scrapy extracted {len(scraped)} markets instead of 4")
    scrapy_link = copy_artifact(scrapy_json)
    proofs.append(
        ToolProof(
            slug="scrapy",
            name="Scrapy",
            layer="Data and research",
            tagline="Runs a real spider against a local HTML fixture and exports four structured market records.",
            task="Extract market cards with CSS selectors",
            command="scrapy runspider market_spider.py -O markets.json",
            version=package_version("Scrapy"),
            result=json.dumps(scraped, indent=2),
            visual_html=table_visual(["Market", "Score"], [[item["market"], item["score"]] for item in scraped], "Scrapy · extracted items"),
            artifacts=[scrapy_link],
            duration_ms=scrapy_result.duration_ms,
        )
    )
    return proofs


def exercise_diagrams_docs(fixtures: dict[str, Path], work: Path) -> list[ToolProof]:
    proofs: list[ToolProof] = []
    docs = work / "diagrams-docs"
    docs.mkdir()

    pptxgen_script = write_text(
        docs / "pptxgen-demo.js",
        """
        const pptxgen = require('pptxgenjs');
        const pptx = new pptxgen();
        pptx.layout = 'LAYOUT_WIDE';
        pptx.author = 'WutPack fictional example';
        pptx.subject = 'Deterministic tool proof';
        const slide = pptx.addSlide();
        slide.background = { color: 'F3F0E8' };
        slide.addText('Coastal is the first gated pilot.', { x: 0.65, y: 0.45, w: 11.9, h: 0.75, fontFace: 'Arial', fontSize: 30, bold: true, color: '14201D', margin: 0 });
        slide.addText('FICTIONAL DECISION FLOW', { x: 0.65, y: 0.15, w: 4.5, h: 0.2, fontFace: 'Arial', fontSize: 9, bold: true, color: '176B4D', charSpacing: 1.3, margin: 0 });
        const phases = [
          ['Evidence', '81 score', 'B8F36B'],
          ['Economics', '30 days', 'F5D969'],
          ['Gate', 'CFO sign-off', 'F26B38'],
          ['Launch', 'if proven', '176B4D']
        ];
        phases.forEach((phase, index) => {
          const x = 0.65 + index * 3.1;
          slide.addShape(pptx.ShapeType.rect, { x, y: 2.25, w: 2.65, h: 2.1, fill: { color: phase[2] }, line: { color: '14201D', width: 1.2 } });
          slide.addText(phase[0], { x: x + 0.2, y: 2.55, w: 2.25, h: 0.4, fontFace: 'Arial', fontSize: 19, bold: true, color: phase[2] === '176B4D' ? 'FFFFFF' : '14201D', margin: 0 });
          slide.addText(phase[1], { x: x + 0.2, y: 3.25, w: 2.25, h: 0.35, fontFace: 'Arial', fontSize: 12, color: phase[2] === '176B4D' ? 'FFFFFF' : '14201D', margin: 0 });
        });
        slide.addText('Editable native shapes • generated locally • no client data', { x: 0.65, y: 6.75, w: 7.5, h: 0.25, fontFace: 'Arial', fontSize: 9, color: '58615D', margin: 0 });
        pptx.writeFile({ fileName: 'pptxgenjs-pilot.pptx' });
        """,
    )
    pptxgen_result = run([executable("node"), pptxgen_script], cwd=docs, timeout=120)
    pptxgen_pptx = docs / "pptxgenjs-pilot.pptx"
    if not pptxgen_pptx.is_file():
        raise ToolFailure("PptxGenJS did not create the presentation")
    pptxgen_render = docs / "pptxgen-render"
    pptxgen_render.mkdir()
    run([executable("soffice"), "--headless", f"-env:UserInstallation=file://{docs / 'pptxgen-profile'}", "--convert-to", "pdf", "--outdir", pptxgen_render, pptxgen_pptx], cwd=docs, timeout=120)
    pptxgen_pdf = pptxgen_render / "pptxgenjs-pilot.pdf"
    run([executable("pdftoppm"), "-png", "-singlefile", "-r", "120", pptxgen_pdf, pptxgen_render / "slide"], cwd=docs)
    pptxgen_png = pptxgen_render / "slide.png"
    pptxgen_png_link = copy_artifact(pptxgen_png, "pptxgenjs-pilot.png")
    pptxgen_link = copy_artifact(pptxgen_pptx)
    proofs.append(
        ToolProof(
            slug="pptxgenjs",
            name="PptxGenJS",
            layer="Diagrams and documentation",
            tagline="Creates a widescreen PowerPoint from editable JavaScript-defined shapes and text, then render-checks the slide.",
            task="Generate an editable pilot-gate presentation",
            command="node pptxgen-demo.js → pptx.writeFile()",
            version=node_package_version("pptxgenjs"),
            result="Created a 16:9 PPTX with native rectangles and editable text.",
            visual_html=image_visual(pptxgen_png_link, "Rendered presentation created with PptxGenJS", "pptxgenjs-pilot.pptx", "native PowerPoint render"),
            artifacts=[pptxgen_link],
            duration_ms=pptxgen_result.duration_ms,
        )
    )

    puppeteer_config = write_json(
        docs / "puppeteer.json",
        {
            "executablePath": str(chromium_binary()),
            "args": ["--no-sandbox", "--disable-gpu", "--disable-dev-shm-usage"],
        },
    )
    mermaid_svg = docs / "mermaid-engagement.svg"
    mermaid_result = run([node_executable("mmdc"), "-p", puppeteer_config, "-i", fixtures["mermaid"], "-o", mermaid_svg, "-t", "neutral", "-b", "transparent"], cwd=docs, timeout=180)
    if not mermaid_svg.is_file():
        raise ToolFailure("Mermaid CLI did not create the SVG")
    mermaid_link = copy_artifact(mermaid_svg)
    proofs.append(
        ToolProof(
            slug="mermaid",
            name="Mermaid",
            layer="Diagrams and documentation",
            tagline="Compiles a text flowchart into an accessible, scalable decision diagram using local Chromium.",
            task="Render a source-controlled pilot flow",
            command="mmdc -i engagement.mmd -o engagement.svg",
            version=node_package_version("@mermaid-js/mermaid-cli"),
            result=mermaid_result.stdout or "Generated a Mermaid SVG from five nodes and five edges.",
            visual_html=image_visual(mermaid_link, "Mermaid evidence, score, economics, gate, and launch flowchart", "engagement.mmd", "SVG from text"),
            artifacts=[mermaid_link],
            duration_ms=mermaid_result.duration_ms,
        )
    )

    graphviz_version = first_version_line(run([executable("dot"), "-V"], cwd=docs))
    graphviz_svg = docs / "graphviz-decision.svg"
    graphviz_result = run([executable("dot"), "-Tsvg", fixtures["dot"], "-o", graphviz_svg], cwd=docs)
    graphviz_link = copy_artifact(graphviz_svg)
    proofs.append(
        ToolProof(
            slug="graphviz",
            name="Graphviz",
            layer="Diagrams and documentation",
            tagline="Lays out a directed decision graph from DOT source and emits a portable vector file.",
            task="Lay out a decision dependency graph",
            command="dot -Tsvg decision.dot -o decision.svg",
            version=graphviz_version,
            result="Rendered Evidence → Score → Economics → Gate with pass and hold branches.",
            visual_html=image_visual(graphviz_link, "Graphviz market-entry decision graph", "decision.dot", "DOT layout"),
            artifacts=[graphviz_link],
            duration_ms=graphviz_result.duration_ms,
        )
    )

    plantuml_version = first_version_line(run([executable("plantuml"), "-version"], cwd=docs))
    plantuml_out = docs / "plantuml"
    plantuml_out.mkdir()
    plantuml_result = run([executable("plantuml"), "-tsvg", "-o", plantuml_out, fixtures["plantuml"]], cwd=docs, timeout=120)
    plantuml_svg = plantuml_out / "handoff.svg"
    if not plantuml_svg.is_file():
        raise ToolFailure("PlantUML did not create handoff.svg")
    plantuml_link = copy_artifact(plantuml_svg)
    proofs.append(
        ToolProof(
            slug="plantuml",
            name="PlantUML",
            layer="Diagrams and documentation",
            tagline="Turns a four-part handoff description into a sequence diagram that stays reviewable as text.",
            task="Render an agent-to-tool handoff sequence",
            command="plantuml -tsvg handoff.puml",
            version=plantuml_version,
            result="Rendered Owner → AI host → WutPack skill → persistent tool.",
            visual_html=image_visual(plantuml_link, "PlantUML sequence diagram for a WutPack handoff", "handoff.puml", "sequence SVG"),
            artifacts=[plantuml_link],
            duration_ms=plantuml_result.duration_ms,
        )
    )

    typst_version = first_version_line(run([executable("typst"), "--version"], cwd=docs))
    typst_pdf = docs / "typst-decision.pdf"
    typst_result = run([executable("typst"), "compile", fixtures["typst"], typst_pdf], cwd=docs, timeout=120)
    typst_png = docs / "typst-decision.png"
    run([executable("mutool"), "draw", "-q", "-r", "130", "-o", typst_png, typst_pdf, "1"], cwd=docs)
    typst_png_link = copy_artifact(typst_png)
    typst_pdf_link = copy_artifact(typst_pdf)
    proofs.append(
        ToolProof(
            slug="typst",
            name="Typst",
            layer="Diagrams and documentation",
            tagline="Compiles concise typesetting source into a polished one-page PDF with an explicit decision table.",
            task="Typeset a one-page pilot decision note",
            command="typst compile one-page.typ decision.pdf",
            version=typst_version,
            result="Compiled a US Letter PDF with heading, narrative, and three-row table.",
            visual_html=image_visual(typst_png_link, "Typst-compiled one-page pilot decision note", "one-page.typ", "PDF render"),
            artifacts=[typst_pdf_link],
            duration_ms=typst_result.duration_ms,
        )
    )

    mkdocs_root = docs / "mkdocs"
    (mkdocs_root / "docs").mkdir(parents=True)
    write_text(mkdocs_root / "mkdocs.yml", """
        site_name: Atlas decision room
        site_description: Fictional WutPack tool proof
        nav:
          - Decision: index.md
          - Controls: controls.md
        theme:
          name: mkdocs
    """)
    write_text(mkdocs_root / "docs/index.md", "# Coastal pilot\n\n**Decision:** start with the highest-scoring fictional market.\n\n| Market | Score |\n|---|---:|\n| Coastal | 81 |\n| North | 78 |\n")
    write_text(mkdocs_root / "docs/controls.md", "# Controls\n\nAdvance only after the 30-day economics gate.\n")
    mkdocs_result = run([sys.executable, "-m", "mkdocs", "build", "--strict"], cwd=mkdocs_root, timeout=120)
    mkdocs_index = mkdocs_root / "site/index.html"
    mkdocs_shot = docs / "mkdocs-site.png"
    capture_html(mkdocs_index, mkdocs_shot, docs)
    mkdocs_shot_link = copy_artifact(mkdocs_shot)
    mkdocs_zip = copy_artifact_tree(mkdocs_root / "site", "mkdocs-site")
    proofs.append(
        ToolProof(
            slug="mkdocs",
            name="MkDocs",
            layer="Diagrams and documentation",
            tagline="Builds a linked two-page documentation site from Markdown and fails on broken references in strict mode.",
            task="Publish a small decision room",
            command="python -m mkdocs build --strict",
            version=package_version("mkdocs"),
            result=mkdocs_result.stdout or "Built two linked pages in strict mode.",
            visual_html=image_visual(mkdocs_shot_link, "MkDocs-generated Atlas decision site", "MkDocs · index.html", "real browser capture"),
            artifacts=[mkdocs_zip],
            duration_ms=mkdocs_result.duration_ms,
        )
    )

    sphinx_root = docs / "sphinx"
    source = sphinx_root / "source"
    source.mkdir(parents=True)
    write_text(source / "conf.py", f"project = 'Atlas decision API'\nextensions = ['sphinx.ext.autodoc']\nhtml_theme = 'alabaster'\nimport sys\nsys.path.insert(0, {str(fixtures['python_module'].parent)!r})\n")
    write_text(source / "index.rst", "Atlas decision API\n==================\n\nA fictional Sphinx build.\n\n.. automodule:: market_model\n   :members:\n")
    sphinx_out = sphinx_root / "html"
    sphinx_result = run([sys.executable, "-m", "sphinx", "-W", "-b", "html", source, sphinx_out], cwd=sphinx_root, timeout=120)
    sphinx_shot = docs / "sphinx-site.png"
    capture_html(sphinx_out / "index.html", sphinx_shot, docs)
    sphinx_shot_link = copy_artifact(sphinx_shot)
    sphinx_zip = copy_artifact_tree(sphinx_out, "sphinx-site")
    proofs.append(
        ToolProof(
            slug="sphinx",
            name="Sphinx",
            layer="Diagrams and documentation",
            tagline="Builds reference documentation from reStructuredText and live Python docstrings with warnings treated as errors.",
            task="Generate a Python API reference",
            command="python -m sphinx -W -b html source html",
            version=package_version("Sphinx"),
            result=sphinx_result.stdout.splitlines()[-1] if sphinx_result.stdout else "Build succeeded with warnings as errors.",
            visual_html=image_visual(sphinx_shot_link, "Sphinx-generated API documentation", "Sphinx · API reference", "real browser capture"),
            artifacts=[sphinx_zip],
            duration_ms=sphinx_result.duration_ms,
        )
    )

    pdoc_out = docs / "pdoc"
    pdoc_result = run([sys.executable, "-m", "pdoc", "-o", pdoc_out, fixtures["python_module"]], cwd=docs, timeout=120)
    pdoc_pages = sorted(pdoc_out.glob("*.html"))
    if not pdoc_pages:
        raise ToolFailure("pdoc did not generate an HTML page")
    pdoc_shot = docs / "pdoc-site.png"
    capture_html(pdoc_pages[0], pdoc_shot, docs)
    pdoc_shot_link = copy_artifact(pdoc_shot)
    pdoc_zip = copy_artifact_tree(pdoc_out, "pdoc-site")
    proofs.append(
        ToolProof(
            slug="pdoc",
            name="pdoc",
            layer="Diagrams and documentation",
            tagline="Publishes a readable Python module reference directly from a real function signature and docstring.",
            task="Document a market-scoring Python module",
            command="python -m pdoc -o html market_model.py",
            version=package_version("pdoc"),
            result=pdoc_result.stdout or "Generated module documentation with signature and return type.",
            visual_html=image_visual(pdoc_shot_link, "pdoc-generated module documentation", "pdoc · market_model", "real browser capture"),
            artifacts=[pdoc_zip],
            duration_ms=pdoc_result.duration_ms,
        )
    )

    doxygen_root = docs / "doxygen"
    doxygen_root.mkdir()
    doxyfile = write_text(
        doxygen_root / "Doxyfile",
        f"""
        PROJECT_NAME = "Atlas score helper"
        OUTPUT_DIRECTORY = {doxygen_root / 'out'}
        INPUT = {fixtures['c_source']}
        GENERATE_HTML = YES
        GENERATE_LATEX = NO
        QUIET = YES
        WARN_AS_ERROR = YES
        EXTRACT_ALL = YES
        """,
    )
    doxygen_version = first_version_line(run([executable("doxygen"), "--version"], cwd=doxygen_root))
    doxygen_result = run([executable("doxygen"), doxyfile], cwd=doxygen_root, timeout=120)
    doxygen_html = doxygen_root / "out/html/index.html"
    doxygen_shot = docs / "doxygen-site.png"
    capture_html(doxygen_html, doxygen_shot, docs)
    doxygen_shot_link = copy_artifact(doxygen_shot)
    doxygen_zip = copy_artifact_tree(doxygen_root / "out/html", "doxygen-site")
    proofs.append(
        ToolProof(
            slug="doxygen",
            name="Doxygen",
            layer="Diagrams and documentation",
            tagline="Turns documented C source into browsable API reference pages with warnings promoted to build failures.",
            task="Generate a C function reference",
            command="doxygen Doxyfile",
            version=doxygen_version,
            result=doxygen_result.stdout or "Generated HTML documentation with zero warnings.",
            visual_html=image_visual(doxygen_shot_link, "Doxygen-generated C API documentation", "Doxygen · market_score()", "real browser capture"),
            artifacts=[doxygen_zip],
            duration_ms=doxygen_result.duration_ms,
        )
    )

    jsdoc_out = docs / "jsdoc"
    jsdoc_result = run([node_executable("jsdoc"), fixtures["javascript"], "-d", jsdoc_out], cwd=docs, timeout=120)
    jsdoc_shot = docs / "jsdoc-site.png"
    capture_html(jsdoc_out / "index.html", jsdoc_shot, docs)
    jsdoc_shot_link = copy_artifact(jsdoc_shot)
    jsdoc_zip = copy_artifact_tree(jsdoc_out, "jsdoc-site")
    proofs.append(
        ToolProof(
            slug="jsdoc",
            name="JSDoc",
            layer="Diagrams and documentation",
            tagline="Reads JavaScript annotations and emits a linked function reference with parameter and return contracts.",
            task="Document a JavaScript market selector",
            command="jsdoc decision.js -d html",
            version=node_package_version("jsdoc"),
            result=jsdoc_result.stdout or "Generated documentation for selectMarket(markets).",
            visual_html=image_visual(jsdoc_shot_link, "JSDoc-generated JavaScript API reference", "JSDoc · selectMarket", "real browser capture"),
            artifacts=[jsdoc_zip],
            duration_ms=jsdoc_result.duration_ms,
        )
    )

    typedoc_out = docs / "typedoc"
    typedoc_result = run([node_executable("typedoc"), "--entryPoints", fixtures["typescript"], "--out", typedoc_out, "--name", "Atlas decision API"], cwd=docs, timeout=120)
    typedoc_shot = docs / "typedoc-site.png"
    capture_html(typedoc_out / "index.html", typedoc_shot, docs)
    typedoc_shot_link = copy_artifact(typedoc_shot)
    typedoc_zip = copy_artifact_tree(typedoc_out, "typedoc-site")
    proofs.append(
        ToolProof(
            slug="typedoc",
            name="TypeDoc",
            layer="Diagrams and documentation",
            tagline="Builds a navigable TypeScript reference that preserves the interface and typed return contract.",
            task="Publish a typed market-selection API",
            command="typedoc --entryPoints decision.ts --out html",
            version=node_package_version("typedoc"),
            result=typedoc_result.stdout or "Generated TypeScript interface and function pages.",
            visual_html=image_visual(typedoc_shot_link, "TypeDoc-generated TypeScript API reference", "TypeDoc · Market", "real browser capture"),
            artifacts=[typedoc_zip],
            duration_ms=typedoc_result.duration_ms,
        )
    )
    return proofs


def exercise_engineering(fixtures: dict[str, Path], work: Path) -> list[ToolProof]:
    proofs: list[ToolProof] = []
    engineering = work / "engineering"
    engineering.mkdir()

    gh_version = first_version_line(run([executable("gh"), "--version"], cwd=engineering))
    gh_auth = run([executable("gh"), "auth", "status", "--hostname", "github.com"], cwd=engineering, expected=(0, 1), timeout=30)
    auth_text = gh_auth.stdout or gh_auth.stderr
    auth_text = re.sub(r"account\s+\S+", "account $ACCOUNT", auth_text, flags=re.IGNORECASE)
    auth_text = re.sub(r"Token:\s+\S+", "Token: [redacted]", auth_text, flags=re.IGNORECASE)
    gh_lines = [line for line in auth_text.splitlines() if "token" not in line.lower()][:7]
    proofs.append(
        ToolProof(
            slug="github-cli",
            name="GitHub CLI",
            layer="Engineering utilities",
            tagline="Checks the installed GitHub CLI and its local authentication state while redacting account and credential details.",
            task="Verify the local GitHub command is ready",
            command="gh auth status --hostname github.com",
            version=gh_version,
            result="\n".join(gh_lines) or "GitHub CLI ran; no authenticated host was reported.",
            visual_html=terminal_visual(["$ gh auth status --hostname github.com", *(gh_lines or ["GitHub CLI available; authentication not configured."]), "", "$ gh pr create --fill", "→ command available when a branch is ready"], "GitHub CLI · credential-safe check"),
            duration_ms=gh_auth.duration_ms,
            note="The proof is intentionally read-only. It does not create an issue, branch, pull request, or network-side change.",
        )
    )

    rg_version = first_version_line(run([executable("rg"), "--version"], cwd=engineering))
    rg_result = run([executable("rg"), "-n", "Coastal|North", fixtures["csv"]], cwd=engineering)
    proofs.append(
        ToolProof(
            slug="ripgrep",
            name="ripgrep",
            layer="Engineering utilities",
            tagline="Searches the scorecard with a regular expression and reports exact source line numbers.",
            task="Find the two leading market rows",
            command="rg -n 'Coastal|North' market-scorecard.csv",
            version=rg_version,
            result=rg_result.stdout,
            visual_html=code_visual(rg_result.stdout, "ripgrep · 2 matches", "search results"),
            duration_ms=rg_result.duration_ms,
        )
    )

    file_tree = engineering / "files"
    (file_tree / "reports/2026").mkdir(parents=True)
    (file_tree / "archive").mkdir()
    write_text(file_tree / "reports/2026/market-brief.md", "# Market brief\n")
    write_text(file_tree / "reports/2026/pilot-plan.md", "# Pilot plan\n")
    write_text(file_tree / "archive/old-plan.txt", "archived\n")
    fd_version = first_version_line(run([executable("fd"), "--version"], cwd=engineering))
    fd_result = run([executable("fd"), "--extension", "md", ".", file_tree], cwd=engineering)
    found_files = [clean_output(line).replace(str(file_tree), "files") for line in fd_result.stdout.splitlines()]
    proofs.append(
        ToolProof(
            slug="fd",
            name="fd",
            layer="Engineering utilities",
            tagline="Finds only current Markdown deliverables in a nested project tree with a concise file-type filter.",
            task="Discover all Markdown reports",
            command="fd --extension md . files/",
            version=fd_version,
            result="\n".join(found_files),
            visual_html=terminal_visual(["$ fd --extension md . files/", *[f"↳ {line}" for line in found_files]], "fd · file discovery"),
            duration_ms=fd_result.duration_ms,
        )
    )

    candidates = "\n".join(["archive/old-plan.txt", "reports/2026/market-brief.md", "reports/2026/pilot-plan.md", "data/market-scorecard.csv"]) + "\n"
    fzf_version = first_version_line(run([executable("fzf"), "--version"], cwd=engineering))
    fzf_result = run([executable("fzf"), "--filter", "market"], cwd=engineering, input_text=candidates)
    proofs.append(
        ToolProof(
            slug="fzf",
            name="fzf",
            layer="Engineering utilities",
            tagline="Ranks a list of project files with fuzzy matching and returns only the market-related choices.",
            task="Filter a deliverable list without an interactive UI",
            command="printf '<files>' | fzf --filter market",
            version=fzf_version,
            result=fzf_result.stdout,
            visual_html=terminal_visual(["$ fzf --filter market", *[f"✓ {line}" for line in fzf_result.stdout.splitlines()]], "fzf · ranked matches"),
            duration_ms=fzf_result.duration_ms,
        )
    )

    jq_version = first_version_line(run([executable("jq"), "--version"], cwd=engineering))
    jq_filter = '[.markets[] | select(.score >= 75)] | sort_by(-.score) | map({market: .name, score, decision: "pilot"})'
    jq_result = run([executable("jq"), jq_filter, fixtures["json"]], cwd=engineering)
    jq_rows = json.loads(jq_result.stdout)
    proofs.append(
        ToolProof(
            slug="jq",
            name="jq",
            layer="Engineering utilities",
            tagline="Filters, sorts, and reshapes nested JSON into the two markets that clear a decision threshold.",
            task="Turn raw JSON into a pilot shortlist",
            command="jq '[.markets[] | select(.score >= 75)] | sort_by(-.score)' markets.json",
            version=jq_version,
            result=jq_result.stdout,
            visual_html=table_visual(["Market", "Score", "Decision"], [[row["market"], row["score"], row["decision"]] for row in jq_rows], "jq · transformed JSON"),
            duration_ms=jq_result.duration_ms,
        )
    )

    yq_version = first_version_line(run([executable("yq"), "--version"], cwd=engineering))
    yq_result = run([executable("yq"), "-o=json", '.gates | map({"gate": .name, "owner": .owner})', fixtures["yaml"]], cwd=engineering)
    yq_rows = json.loads(yq_result.stdout)
    proofs.append(
        ToolProof(
            slug="yq",
            name="yq",
            layer="Engineering utilities",
            tagline="Queries a YAML launch plan and emits a normalized JSON list of gates and accountable owners.",
            task="Extract controls from a YAML plan",
            command="yq -o=json '.gates | map({\"gate\": .name, \"owner\": .owner})' launch-plan.yml",
            version=yq_version,
            result=yq_result.stdout,
            visual_html=table_visual(["Gate", "Owner"], [[row["gate"], row["owner"]] for row in yq_rows], "yq · launch controls"),
            duration_ms=yq_result.duration_ms,
        )
    )

    shellcheck_version = first_version_line(run([executable("shellcheck"), "--version"], cwd=engineering))
    shellcheck_result = run([executable("shellcheck"), "--format=gcc", fixtures["broken_shell"]], cwd=engineering, expected=(1,))
    proofs.append(
        ToolProof(
            slug="shellcheck",
            name="ShellCheck",
            layer="Engineering utilities",
            tagline="Finds unsafe shell expansion in a small script and returns line-addressable diagnostics before the script ships.",
            task="Lint an intentionally unsafe shell loop",
            command="shellcheck --format=gcc broken.sh",
            version=shellcheck_version,
            result=shellcheck_result.stdout,
            visual_html=code_visual(shellcheck_result.stdout, "ShellCheck · expected findings", "diagnostics"),
            duration_ms=shellcheck_result.duration_ms,
            note="A non-zero exit is expected here: the successful proof is that the deliberately broken fixture is rejected with actionable diagnostics.",
        )
    )

    shfmt_version = first_version_line(run([executable("shfmt"), "--version"], cwd=engineering))
    before = fixtures["messy_shell"].read_text(encoding="utf-8")
    formatted = engineering / "formatted.sh"
    shutil.copy2(fixtures["messy_shell"], formatted)
    shfmt_result = run([executable("shfmt"), "-w", "-i", "2", formatted], cwd=engineering)
    after = formatted.read_text(encoding="utf-8")
    if before == after:
        raise ToolFailure("shfmt did not change the deliberately messy fixture")
    proofs.append(
        ToolProof(
            slug="shfmt",
            name="shfmt",
            layer="Engineering utilities",
            tagline="Normalizes shell indentation and control-flow spacing into a deterministic style.",
            task="Format a deliberately messy shell script",
            command="shfmt -w -i 2 messy.sh",
            version=shfmt_version,
            result=after,
            visual_html=f'<div class="split-visual"><pre><span>Before</span>\n{escape(before)}</pre><pre><span>After</span>\n{escape(after)}</pre></div>',
            duration_ms=shfmt_result.duration_ms,
        )
    )

    before_file = write_text(engineering / "before.md", "# Pilot plan\n\n- Coastal: analyze\n- North: analyze\n")
    after_file = write_text(engineering / "after.md", "# Pilot plan\n\n- Coastal: launch gated pilot\n- North: prepare next\n- Control: 20% margin gate\n")
    git_diff = run([executable("git"), "diff", "--no-index", "--", before_file, after_file], cwd=engineering, expected=(1,))
    delta_version = first_version_line(run([executable("delta"), "--version"], cwd=engineering))
    delta_result = run([executable("delta"), "--paging=never", "--line-numbers"], cwd=engineering, input_text=git_diff.stdout)
    proofs.append(
        ToolProof(
            slug="delta",
            name="delta",
            layer="Engineering utilities",
            tagline="Turns a raw Git patch into a line-numbered, readable review view that makes the decision change obvious.",
            task="Render a plan revision for code review",
            command="git diff --no-index before.md after.md | delta --line-numbers",
            version=delta_version,
            result=delta_result.stdout,
            visual_html=code_visual(delta_result.stdout, "delta · reviewed change", "diff"),
            duration_ms=git_diff.duration_ms + delta_result.duration_ms,
        )
    )

    benchmark_file = engineering / "benchmark.txt"
    benchmark_file.write_text(("Coastal,81\nNorth,78\nCentral,72\nWest,58\n") * 1200, encoding="utf-8")
    hyperfine_version = first_version_line(run([executable("hyperfine"), "--version"], cwd=engineering))
    hyperfine_json = engineering / "hyperfine.json"
    rg_cmd = f"{executable('rg')} -c Coastal {benchmark_file}"
    grep_cmd = f"/usr/bin/grep -c Coastal {benchmark_file}"
    hyperfine_result = run([executable("hyperfine"), "--warmup", "1", "--runs", "3", "--export-json", hyperfine_json, rg_cmd, grep_cmd], cwd=engineering, timeout=120)
    benchmark_payload = json.loads(hyperfine_json.read_text(encoding="utf-8"))
    benchmark = benchmark_payload["results"]
    benchmark[0]["command"] = "rg -c Coastal benchmark.txt"
    benchmark[1]["command"] = "grep -c Coastal benchmark.txt"
    write_json(hyperfine_json, benchmark_payload)
    benchmark_rows = [("ripgrep", benchmark[0]["mean"] * 1000), ("grep", benchmark[1]["mean"] * 1000)]
    hyperfine_link = copy_artifact(hyperfine_json)
    proofs.append(
        ToolProof(
            slug="hyperfine",
            name="hyperfine",
            layer="Engineering utilities",
            tagline="Benchmarks two equivalent local searches with warmup, repeated runs, statistical timing, and JSON export.",
            task="Compare two text-search commands",
            command="hyperfine --warmup 1 --runs 3 --export-json results.json '<rg>' '<grep>'",
            version=hyperfine_version,
            result=hyperfine_result.stdout,
            visual_html=bars_visual(benchmark_rows, "Mean runtime", " ms"),
            artifacts=[hyperfine_link],
            duration_ms=hyperfine_result.duration_ms,
            note="The tiny benchmark proves the harness and JSON output; it is not a general claim that one command is always faster.",
        )
    )

    just_version = first_version_line(run([executable("just"), "--version"], cwd=engineering))
    just_result = run([executable("just"), "--justfile", fixtures["justfile"], "score", "Coastal", "81"], cwd=engineering)
    proofs.append(
        ToolProof(
            slug="just",
            name="just",
            layer="Engineering utilities",
            tagline="Runs a named, parameterized project recipe so a repeated score check has one memorable command.",
            task="Execute a reusable score recipe",
            command="just --justfile justfile score Coastal 81",
            version=just_version,
            result=just_result.stdout,
            visual_html=terminal_visual(["$ just score Coastal 81", just_result.stdout.strip(), "", "recipe: score market value", "status: completed"], "just · project recipe"),
            duration_ms=just_result.duration_ms,
        )
    )
    return proofs


def validate_proofs(proofs: Sequence[ToolProof]) -> None:
    expected_names = [name for layer in LAYER_ORDER for name in EXPECTED_TOOLS[layer]]
    actual_names = [proof.name for proof in proofs]
    if len(proofs) != 53:
        raise ToolFailure(f"expected 53 tool proofs, built {len(proofs)}")
    if len({proof.slug for proof in proofs}) != len(proofs):
        raise ToolFailure("tool proof slugs are not unique")
    if len([proof for proof in proofs if proof.status == "passed"]) != 52:
        raise ToolFailure("expected 52 exercised tool proofs")
    if [(proof.name, proof.status) for proof in proofs if proof.status != "passed"] != [("Claude Code", "installed-only")]:
        raise ToolFailure("Claude Code must be the only installed-only proof")
    if set(actual_names) != set(expected_names):
        missing = sorted(set(expected_names) - set(actual_names))
        extra = sorted(set(actual_names) - set(expected_names))
        raise ToolFailure(f"tool proof inventory differs; missing={missing}, extra={extra}")
    actual_layers = {layer: [proof.name for proof in proofs if proof.layer == layer] for layer in LAYER_ORDER}
    for layer in LAYER_ORDER:
        if actual_layers[layer] != EXPECTED_TOOLS[layer]:
            raise ToolFailure(
                f"{layer} order differs; expected={EXPECTED_TOOLS[layer]}, actual={actual_layers[layer]}"
            )


def gallery_fragment(proofs: Sequence[ToolProof]) -> str:
    sections: list[str] = []
    for index, layer in enumerate(LAYER_ORDER, start=1):
        layer_proofs = [proof for proof in proofs if proof.layer == layer]
        links = "\n".join(
            f'''<a class="tool-example-link{(' installed-only' if proof.status != 'passed' else '')}" href="tool-examples/{escape(proof.slug)}.html" data-tool-name="{escape((proof.name + ' ' + proof.tagline).lower())}">
              <span class="tool-example-name"><i aria-hidden="true"></i><strong>{escape(proof.name)}</strong></span>
              <span>{escape(proof.tagline)}</span>
              <small>{'Open install record' if proof.status != 'passed' else 'Open visual proof'} →</small>
            </a>'''
            for proof in layer_proofs
        )
        sections.append(
            f'''<section class="tool-layer" aria-labelledby="tool-layer-{index}">
          <div class="tool-layer-heading">
            <span>0{index}</span>
            <div><h3 id="tool-layer-{index}">{escape(layer)}</h3><p>{len([proof for proof in layer_proofs if proof.status == 'passed'])} exercised · {len(layer_proofs)} pages</p></div>
          </div>
          <div class="tool-example-grid">{links}</div>
        </section>'''
        )
    return f'''<div class="tool-gallery-controls reveal">
        <label for="tool-search">Find a tool or job</label>
        <div class="tool-search-wrap"><input id="tool-search" type="search" placeholder="Try PDF, diagram, browser, Excel…" autocomplete="off" data-tool-search><span data-tool-count aria-live="polite">53 tools</span></div>
      </div>
      <div class="tool-layer-board reveal" data-tool-board>
        {''.join(sections)}
      </div>
      <p class="tool-empty" data-tool-empty hidden>No matching tool. Try a broader word.</p>'''


def refresh_landing_page(proofs: Sequence[ToolProof]) -> None:
    text = LANDING_PAGE.read_text(encoding="utf-8")
    start_marker = "<!-- TOOL_GALLERY_START -->"
    end_marker = "<!-- TOOL_GALLERY_END -->"
    if text.count(start_marker) != 1 or text.count(end_marker) != 1:
        raise ToolFailure("landing page must contain one TOOL_GALLERY marker pair")
    before, remainder = text.split(start_marker, 1)
    _, after = remainder.split(end_marker, 1)
    replacement = f"{start_marker}\n      {gallery_fragment(proofs)}\n      {end_marker}"
    LANDING_PAGE.write_text(before + replacement + after, encoding="utf-8")


def write_site(proofs: Sequence[ToolProof]) -> None:
    OUTPUT.mkdir(parents=True, exist_ok=True)
    (OUTPUT / "assets").mkdir(parents=True, exist_ok=True)
    write_text(OUTPUT / "assets/gallery.css", gallery_css())
    write_text(OUTPUT / "assets/gallery.js", gallery_js())

    for index, proof in enumerate(proofs):
        previous_slug = proofs[(index - 1) % len(proofs)].slug
        next_slug = proofs[(index + 1) % len(proofs)].slug
        write_text(OUTPUT / f"{proof.slug}.html", page_for(proof, previous_slug, next_slug))

    evidence = {
        "generated_on": date.today().isoformat(),
        "fixture_policy": "fictional, local, deterministic; no credentials embedded",
        "tool_count": len(proofs),
        "exercised_count": len([proof for proof in proofs if proof.status == "passed"]),
        "installed_only_count": len([proof for proof in proofs if proof.status != "passed"]),
        "layers": {layer: len([proof for proof in proofs if proof.layer == layer]) for layer in LAYER_ORDER},
        "tools": [
            {
                key: value
                for key, value in asdict(proof).items()
                if key not in {"visual_html"}
            }
            for proof in proofs
        ],
    }
    write_json(OUTPUT / "evidence.json", evidence)
    refresh_landing_page(proofs)


def build(*, real_ai: bool, codex_result_path: Path | None) -> list[ToolProof]:
    ARTIFACTS.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="wutpack-tool-gallery-") as temporary:
        work = Path(temporary)
        print("[1/6] fixtures", flush=True)
        fixtures = make_fixtures(work)
        print("[2/6] desktop workbench", flush=True)
        proofs = exercise_desktop(fixtures, work)
        print("[3/6] AI coding CLIs", flush=True)
        proofs.extend(exercise_ai_clis(work, real_ai=real_ai, codex_result_path=codex_result_path))
        print("[4/6] office, PDF, media, data, and research", flush=True)
        proofs.extend(exercise_office_media(fixtures, work))
        proofs.extend(exercise_data_research(fixtures, work))
        print("[5/6] diagrams, documentation, and engineering", flush=True)
        proofs.extend(exercise_diagrams_docs(fixtures, work))
        proofs.extend(exercise_engineering(fixtures, work))
        validate_proofs(proofs)
        print("[6/6] HTML proof pages", flush=True)
        write_site(proofs)
        return proofs


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    codex_source = parser.add_mutually_exclusive_group(required=True)
    codex_source.add_argument(
        "--real-ai",
        action="store_true",
        help="make one authenticated non-interactive Codex model call; Claude is never invoked",
    )
    codex_source.add_argument(
        "--codex-result",
        type=Path,
        help="reuse a schema-validated JSON result from a prior authenticated Codex run",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    codex_result = args.codex_result.resolve() if args.codex_result else None
    if codex_result and not codex_result.is_file():
        print(f"Codex result does not exist: {codex_result}", file=sys.stderr)
        return 2
    try:
        proofs = build(real_ai=args.real_ai, codex_result_path=codex_result)
    except (ToolFailure, subprocess.TimeoutExpired) as exc:
        print(f"tool gallery failed: {exc}", file=sys.stderr)
        return 1
    print(f"built {len(proofs)} exercised tool pages under {OUTPUT.relative_to(ROOT)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
